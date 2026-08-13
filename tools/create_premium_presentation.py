"""
Premium IBM Presentation Generator Tool
Carbon Design System — polished, pixel-perfect IBM-branded slides.

The entire generation engine is inlined here so there are zero external file
dependencies at WXO runtime — identical pattern to how create_presentation.py
inlines pptx_workflow via a local sys.path insert.

Themes : carbon_light | carbon_dark | consulting | professional | teal | executive
Slides : title | content | section | two_column | stats | quote | thank_you

Content fields that accept bullet lists (content, left, right) accept EITHER:
  • A JSON array string:  '["Point one", "Point two"]'
  • Plain newline / dash-delimited text: "- Point one\n- Point two"

Stats field accepts EITHER:
  • A JSON array string: '[{"number": "85%", "label": "Migrated"}]'
  • Plain newline text:  "85% – Workloads migrated\n40% – Faster processing"
"""

import json
import os
import re
import sys
import tempfile
from typing import Dict, List, Optional

from pydantic import BaseModel, Field
from ibm_watsonx_orchestrate.agent_builder.tools import tool

# python-pptx — available via requirements.txt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ══════════════════════════════════════════════════════════════════════════════
# IBM Carbon palette
# ══════════════════════════════════════════════════════════════════════════════

_C = {
    'ibm_blue_60': '0F62FE', 'ibm_blue_70': '0043CE', 'ibm_blue_80': '002D9C',
    'ibm_blue_40': '4589FF', 'ibm_blue_20': 'D0E2FF', 'ibm_blue_10': 'EDF5FF',
    'gray_100': '161616', 'gray_90': '262626', 'gray_80': '393939',
    'gray_70': '525252', 'gray_60': '6F6F6F', 'gray_40': 'A8A8A8',
    'gray_30': 'C6C6C6', 'gray_20': 'E0E0E0', 'gray_10': 'F4F4F4',
    'cyan_50': '1192E8', 'cyan_70': '00539A', 'cyan_20': 'BAE6FF',
    'teal_50': '009D9A', 'teal_70': '004144', 'teal_20': '9EF0F0',
    'purple_50': '8A3FFC', 'purple_70': '491D8B', 'purple_20': 'E8DAFF',
    'magenta_50': 'EE5396', 'magenta_70': '9F1853',
    'green_50': '24A148', 'green_70': '044317',
    'white': 'FFFFFF', 'black': '000000',
    'warm_gray_10': 'F7F3F2', 'warm_gray_80': '4D3D35',
    'red_60': 'DA1E28', 'red_80': '750E13',
    'orange_40': 'FF832B',
}

_FONT = 'IBM Plex Sans'

_THEMES = {
    'carbon_light': dict(
        bg=_C['white'],       panel=_C['ibm_blue_60'],  panel2=_C['ibm_blue_70'],
        accent=_C['ibm_blue_60'], title_fg=_C['gray_100'], body_fg=_C['gray_70'],
        rule=_C['gray_20'],   footer_fg=_C['gray_60'],  tag_bg=_C['ibm_blue_10'],
        tag_fg=_C['ibm_blue_60'], bullet=_C['ibm_blue_60'], stat=_C['ibm_blue_60'],
        cover_bg=_C['ibm_blue_60'], cover_stripe=_C['ibm_blue_80'],
    ),
    'carbon_dark': dict(
        bg=_C['gray_100'],    panel=_C['gray_90'],       panel2=_C['gray_80'],
        accent=_C['ibm_blue_40'], title_fg=_C['white'],  body_fg=_C['gray_30'],
        rule=_C['gray_80'],   footer_fg=_C['gray_60'],   tag_bg=_C['gray_80'],
        tag_fg=_C['ibm_blue_40'], bullet=_C['ibm_blue_40'], stat=_C['ibm_blue_40'],
        cover_bg=_C['gray_90'], cover_stripe=_C['ibm_blue_60'],
    ),
    'consulting': dict(
        bg=_C['white'],       panel=_C['ibm_blue_80'],   panel2=_C['ibm_blue_70'],
        accent=_C['cyan_50'], title_fg=_C['gray_100'],   body_fg=_C['gray_70'],
        rule=_C['gray_20'],   footer_fg=_C['gray_60'],   tag_bg=_C['cyan_20'],
        tag_fg=_C['cyan_70'], bullet=_C['ibm_blue_80'],  stat=_C['ibm_blue_80'],
        cover_bg=_C['ibm_blue_80'], cover_stripe=_C['cyan_50'],
    ),
    'professional': dict(
        bg=_C['gray_10'],     panel=_C['purple_50'],     panel2=_C['purple_70'],
        accent=_C['magenta_50'], title_fg=_C['gray_100'], body_fg=_C['gray_70'],
        rule=_C['gray_20'],   footer_fg=_C['gray_60'],   tag_bg=_C['purple_20'],
        tag_fg=_C['purple_70'], bullet=_C['purple_50'],  stat=_C['purple_50'],
        cover_bg=_C['purple_70'], cover_stripe=_C['magenta_50'],
    ),
    'teal': dict(
        bg=_C['white'],       panel=_C['teal_50'],       panel2=_C['teal_70'],
        accent=_C['teal_50'], title_fg=_C['gray_100'],   body_fg=_C['gray_70'],
        rule=_C['gray_20'],   footer_fg=_C['gray_60'],   tag_bg=_C['teal_20'],
        tag_fg=_C['teal_70'], bullet=_C['teal_50'],      stat=_C['teal_50'],
        cover_bg=_C['teal_70'], cover_stripe=_C['teal_50'],
    ),
    'executive': dict(
        bg=_C['warm_gray_10'], panel=_C['gray_100'],     panel2=_C['gray_90'],
        accent=_C['orange_40'], title_fg=_C['gray_100'], body_fg=_C['gray_70'],
        rule=_C['gray_30'],   footer_fg=_C['gray_60'],   tag_bg=_C['gray_20'],
        tag_fg=_C['gray_80'], bullet=_C['gray_100'],     stat=_C['orange_40'],
        cover_bg=_C['gray_100'], cover_stripe=_C['orange_40'],
    ),
}


def _rgb(h: str) -> RGBColor:
    return RGBColor.from_string(h.lstrip('#'))


# ══════════════════════════════════════════════════════════════════════════════
# Generator class (inlined from ibm_presentation_generator.py)
# ══════════════════════════════════════════════════════════════════════════════

class _IBMGen:
    W = Inches(13.33)
    H = Inches(7.5)

    def __init__(self, theme: str = 'carbon_light'):
        self.prs = Presentation()
        self.prs.slide_width  = self.W
        self.prs.slide_height = self.H
        self.tc = _THEMES.get(theme, _THEMES['carbon_light'])

    def _blank_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        for ph in slide.placeholders:
            ph._element.getparent().remove(ph._element)
        return slide

    def _rect(self, slide, x, y, w, h, fill):
        s = slide.shapes.add_shape(1, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
        s.line.fill.background()
        return s

    def _tb(self, slide, x, y, w, h, text, pt, bold=False,
            color='161616', align=PP_ALIGN.LEFT, wrap=True, italic=False, spacing=0):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tb.fill.background()
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.name   = _FONT
        p.font.size   = Pt(pt)
        p.font.bold   = bold
        p.font.italic = italic
        p.font.color.rgb = _rgb(color)
        if spacing:
            p.space_after = Pt(spacing)
        return tb

    def _footer(self, slide, num: int = None, light=False):
        tc = self.tc
        fg = _C['white'] if light else tc['footer_fg']
        self._rect(slide, Inches(0.4), Inches(6.88), Inches(12.53), Inches(0.02),
                   _C['white'] if light else tc['rule'])
        self._tb(slide, Inches(12.1), Inches(6.9), Inches(1.1), Inches(0.4),
                 'IBM', 16, bold=True, color=fg, align=PP_ALIGN.RIGHT)
        if num:
            self._tb(slide, Inches(0.4), Inches(6.9), Inches(0.8), Inches(0.4),
                     str(num), 10, color=fg)

    # ── Title slide ──────────────────────────────────────────────────────────

    def add_title_slide(self, title, subtitle='', date='', num=1):
        slide = self._blank_slide()
        tc = self.tc
        self._rect(slide, 0, 0, self.W, self.H, tc['cover_bg'])
        self._rect(slide, 0, 0, Inches(1.33), self.H, tc['cover_stripe'])
        self._rect(slide, 0, Inches(6.7), self.W, Inches(0.8), tc['panel2'])
        self._rect(slide, Inches(1.6), Inches(4.8), Inches(9.5), Inches(0.055), _C['white'])
        self._tb(slide, Inches(1.6), Inches(1.8), Inches(10.5), Inches(2.7),
                 title, 56, bold=True, color=_C['white'], align=PP_ALIGN.LEFT, wrap=True)
        if subtitle:
            self._tb(slide, Inches(1.6), Inches(5.05), Inches(10.0), Inches(0.9),
                     subtitle, 22, color=_C['white'], align=PP_ALIGN.LEFT)
        if date:
            self._tb(slide, Inches(1.6), Inches(6.15), Inches(6.0), Inches(0.5),
                     date, 14, color=_C['white'], italic=True)
        self._tb(slide, Inches(11.5), Inches(6.85), Inches(1.6), Inches(0.5),
                 'IBM', 20, bold=True, color=_C['white'], align=PP_ALIGN.RIGHT)
        return slide

    # ── Content slide ────────────────────────────────────────────────────────

    def add_content_slide(self, title, content, eyebrow='', layout='bullets', num=2):
        slide = self._blank_slide()
        tc = self.tc
        self._rect(slide, 0, 0, self.W, self.H, tc['bg'])
        self._rect(slide, 0, 0, self.W, Inches(1.45), tc['panel'])
        self._rect(slide, 0, 0, Inches(0.22), Inches(1.45), tc['panel2'])
        self._rect(slide, 0, Inches(1.45), self.W, Inches(0.07), tc['accent'])
        if eyebrow:
            self._tb(slide, Inches(0.45), Inches(0.15), Inches(11.9), Inches(0.35),
                     eyebrow.upper(), 11, color=_C['white'], italic=True)
        self._tb(slide, Inches(0.45), Inches(0.38), Inches(11.9), Inches(1.0),
                 title, 34, bold=True, color=_C['white'], align=PP_ALIGN.LEFT)
        body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(12.1), Inches(4.9))
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(content[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if layout == 'numbered':
                p.text = f"{i + 1}.  {item}"
                p.font.name  = _FONT
                p.font.size  = Pt(20)
                p.font.color.rgb = _rgb(tc['bullet'])
            else:
                r1 = p.add_run()
                r1.text = '\u2014  '
                r1.font.name  = _FONT
                r1.font.size  = Pt(20)
                r1.font.bold  = True
                r1.font.color.rgb = _rgb(tc['bullet'])
                r2 = p.add_run()
                r2.text = item
                r2.font.name  = _FONT
                r2.font.size  = Pt(20)
                r2.font.color.rgb = _rgb(tc['body_fg'])
            p.space_before = Pt(4)
            p.space_after  = Pt(10)
        self._footer(slide, num)
        return slide

    # ── Section divider ──────────────────────────────────────────────────────

    def add_section_divider(self, title, eyebrow='', num=2):
        slide = self._blank_slide()
        tc = self.tc
        self._rect(slide, 0, 0, self.W, self.H, tc['cover_bg'])
        self._rect(slide, 0, 0, Inches(1.33), self.H, tc['cover_stripe'])
        self._rect(slide, 0, Inches(6.5), self.W, Inches(1.0), tc['panel2'])
        label = eyebrow if eyebrow else 'SECTION'
        self._tb(slide, Inches(1.65), Inches(1.4), Inches(8.0), Inches(0.5),
                 label.upper(), 12, color=_C['white'], italic=True)
        self._rect(slide, Inches(1.65), Inches(1.95), Inches(1.2), Inches(0.055), tc['cover_stripe'])
        self._tb(slide, Inches(1.65), Inches(2.1), Inches(9.2), Inches(3.0),
                 title, 52, bold=True, color=_C['white'], align=PP_ALIGN.LEFT, wrap=True)
        self._footer(slide, light=True)
        return slide

    # ── Two-column ───────────────────────────────────────────────────────────

    def add_two_column_slide(self, title, left_content, right_content,
                             left_label='', right_label='', num=2):
        slide = self._blank_slide()
        tc = self.tc
        self._rect(slide, 0, 0, self.W, self.H, tc['bg'])
        self._rect(slide, 0, 0, self.W, Inches(1.45), tc['panel'])
        self._rect(slide, 0, 0, Inches(0.22), Inches(1.45), tc['panel2'])
        self._rect(slide, 0, Inches(1.45), self.W, Inches(0.07), tc['accent'])
        self._tb(slide, Inches(0.45), Inches(0.2), Inches(11.9), Inches(1.05),
                 title, 34, bold=True, color=_C['white'])
        self._rect(slide, Inches(6.61), Inches(1.65), Inches(0.04), Inches(5.0), tc['rule'])
        col_top = Inches(1.65)
        for items, label, x in [
            (left_content,  left_label,  Inches(0.55)),
            (right_content, right_label, Inches(6.85)),
        ]:
            cur_top = col_top
            if label:
                self._rect(slide, x, cur_top, Inches(5.7), Inches(0.42), tc['tag_bg'])
                self._tb(slide, x + Inches(0.12), cur_top + Inches(0.04),
                         Inches(5.4), Inches(0.34), label, 14, bold=True, color=tc['tag_fg'])
                cur_top += Inches(0.55)
            cb = slide.shapes.add_textbox(x, cur_top, Inches(5.7), Inches(4.5))
            cf = cb.text_frame
            cf.word_wrap = True
            for i, item in enumerate(items[:6]):
                p = cf.paragraphs[0] if i == 0 else cf.add_paragraph()
                r1 = p.add_run()
                r1.text = '\u2014  '
                r1.font.name = _FONT
                r1.font.size = Pt(18)
                r1.font.bold = True
                r1.font.color.rgb = _rgb(tc['bullet'])
                r2 = p.add_run()
                r2.text = item
                r2.font.name = _FONT
                r2.font.size = Pt(18)
                r2.font.color.rgb = _rgb(tc['body_fg'])
                p.space_after = Pt(9)
        self._footer(slide, num)
        return slide

    # ── Stats slide ──────────────────────────────────────────────────────────

    def add_stats_slide(self, title, stats, num=2):
        slide = self._blank_slide()
        tc = self.tc
        self._rect(slide, 0, 0, self.W, self.H, tc['bg'])
        self._rect(slide, 0, 0, self.W, Inches(1.45), tc['panel'])
        self._rect(slide, 0, 0, Inches(0.22), Inches(1.45), tc['panel2'])
        self._rect(slide, 0, Inches(1.45), self.W, Inches(0.07), tc['accent'])
        self._tb(slide, Inches(0.45), Inches(0.2), Inches(11.9), Inches(1.05),
                 title, 34, bold=True, color=_C['white'])
        n = min(len(stats), 4)
        if n == 0:
            self._footer(slide, num)
            return slide
        col_w = 12.53 / n
        for i, stat in enumerate(stats[:4]):
            x = Inches(0.4 + i * col_w)
            w = Inches(col_w)
            if i > 0:
                self._rect(slide, x - Inches(0.02), Inches(1.8), Inches(0.04), Inches(4.6), tc['rule'])
            self._tb(slide, x + Inches(0.15), Inches(2.1), w - Inches(0.3), Inches(2.2),
                     stat.get('number', '—'), 72, bold=True, color=tc['stat'], align=PP_ALIGN.LEFT)
            self._rect(slide, x + Inches(0.15), Inches(3.9),
                       Inches(col_w * 0.35), Inches(0.055), tc['accent'])
            self._tb(slide, x + Inches(0.15), Inches(4.1), w - Inches(0.3), Inches(2.3),
                     stat.get('label', ''), 18, color=tc['body_fg'], wrap=True)
            if stat.get('source'):
                self._tb(slide, x + Inches(0.15), Inches(6.2), w - Inches(0.3), Inches(0.5),
                         stat['source'], 11, italic=True, color=tc['footer_fg'])
        self._footer(slide, num)
        return slide

    # ── Quote slide ──────────────────────────────────────────────────────────

    def add_quote_slide(self, title, quote, attribution='', num=2):
        slide = self._blank_slide()
        tc = self.tc
        self._rect(slide, 0, 0, self.W, self.H, tc['bg'])
        self._rect(slide, 0, 0, Inches(0.45), self.H, tc['panel'])
        self._tb(slide, Inches(0.7), Inches(0.15), Inches(11.5), Inches(0.6),
                 title, 22, bold=True, color=tc['accent'])
        self._rect(slide, Inches(0.7), Inches(0.78), Inches(1.5), Inches(0.055), tc['accent'])
        self._tb(slide, Inches(0.7), Inches(0.95), Inches(2.0), Inches(2.0),
                 '\u201C', 120, bold=True, color=tc['accent'])
        self._tb(slide, Inches(0.8), Inches(2.1), Inches(11.3), Inches(3.6),
                 quote, 26, bold=False, color=tc['title_fg'],
                 align=PP_ALIGN.LEFT, wrap=True, italic=True)
        if attribution:
            self._rect(slide, Inches(0.7), Inches(5.95), Inches(0.4), Inches(0.55), tc['accent'])
            self._tb(slide, Inches(1.25), Inches(5.95), Inches(11.0), Inches(0.55),
                     attribution, 16, bold=True, color=tc['title_fg'])
        self._footer(slide, num)
        return slide

    # ── Thank you ────────────────────────────────────────────────────────────

    def add_thank_you_slide(self, message='Thank You', contact='', num=2):
        slide = self._blank_slide()
        tc = self.tc
        self._rect(slide, 0, 0, self.W, self.H, tc['cover_bg'])
        self._rect(slide, 0, 0, Inches(1.33), self.H, tc['cover_stripe'])
        self._rect(slide, 0, Inches(6.5), self.W, Inches(1.0), tc['panel2'])
        self._rect(slide, Inches(1.6), Inches(5.0), Inches(9.5), Inches(0.055), _C['white'])
        self._tb(slide, Inches(1.6), Inches(2.2), Inches(10.5), Inches(2.5),
                 message, 60, bold=True, color=_C['white'], wrap=True)
        if contact:
            self._tb(slide, Inches(1.6), Inches(5.2), Inches(10.0), Inches(0.9),
                     contact, 20, color=_C['white'])
        self._tb(slide, Inches(11.5), Inches(6.85), Inches(1.6), Inches(0.5),
                 'IBM', 20, bold=True, color=_C['white'], align=PP_ALIGN.RIGHT)
        return slide

    def save(self, path: str) -> None:
        self.prs.save(path)


# ══════════════════════════════════════════════════════════════════════════════
# Input parsing helpers
# ══════════════════════════════════════════════════════════════════════════════

def _parse_list(value: Optional[str]) -> List[str]:
    """Parse bullet list — JSON array OR plain newline/dash text."""
    if not value:
        return []
    s = value.strip()
    if s.startswith('['):
        try:
            parsed = json.loads(s)
            if isinstance(parsed, list):
                return [str(x) for x in parsed]
        except json.JSONDecodeError:
            pass
    lines = []
    for line in s.splitlines():
        line = line.strip()
        if not line:
            continue
        line = re.sub(r'^[\-\•\*\u2013\u2014]\s*', '', line)
        line = re.sub(r'^\d+[\.\)]\s*', '', line)
        if line:
            lines.append(line)
    return lines


def _parse_stats(value: Optional[str]) -> List[Dict]:
    """Parse stats — JSON array OR "number – label" lines."""
    if not value:
        return []
    s = value.strip()
    if s.startswith('['):
        try:
            parsed = json.loads(s)
            if isinstance(parsed, list):
                return parsed
        except json.JSONDecodeError:
            pass
    stats = []
    for line in s.splitlines():
        line = line.strip()
        if not line:
            continue
        m = re.split(r'\s*[\u2013\u2014\-\|:]\s*', line, maxsplit=1)
        if len(m) == 2:
            stats.append({'number': m[0].strip(), 'label': m[1].strip()})
        else:
            stats.append({'number': '', 'label': line})
    return stats


# ══════════════════════════════════════════════════════════════════════════════
# Pydantic models
# ══════════════════════════════════════════════════════════════════════════════

class PremiumSlide(BaseModel):
    type: str = Field(
        description=(
            "Slide type. One of: "
            "'title' | 'content' | 'section' | 'two_column' | 'stats' | 'quote' | 'thank_you'"
        )
    )
    title:       Optional[str] = Field(default=None, description="Slide title")
    subtitle:    Optional[str] = Field(default=None, description="Subtitle — title slide only")
    date:        Optional[str] = Field(default=None, description="Date string — title slide, e.g. 'July 22, 2026'")
    eyebrow:     Optional[str] = Field(default=None, description="Short label above heading, e.g. 'The opportunity'")
    content:     Optional[str] = Field(
        default=None,
        description=(
            "Bullet points for 'content' slides. "
            "JSON array '[\"Point one\", \"Point two\"]' OR plain newline-separated text."
        )
    )
    layout:      Optional[str] = Field(default='bullets', description="'bullets' (default) or 'numbered'")
    left:        Optional[str] = Field(default=None, description="Left-column bullets for 'two_column'. Same format as content.")
    right:       Optional[str] = Field(default=None, description="Right-column bullets for 'two_column'. Same format as content.")
    left_label:  Optional[str] = Field(default=None, description="Left column header (two_column slide)")
    right_label: Optional[str] = Field(default=None, description="Right column header (two_column slide)")
    stats:       Optional[str] = Field(
        default=None,
        description=(
            "Stats for 'stats' slides. "
            "JSON array '[{\"number\": \"85%\", \"label\": \"Migrated\"}]' OR "
            "plain text '£2.3M – Annual saving\\n40% – Faster processing'. 2–4 stats."
        )
    )
    quote:       Optional[str] = Field(default=None, description="Pull-quote text (quote slide)")
    attribution: Optional[str] = Field(default=None, description="Quote attribution (quote slide)")
    message:     Optional[str] = Field(default=None, description="Closing message (thank_you slide)")
    contact:     Optional[str] = Field(default=None, description="Contact info (thank_you slide)")


class PremiumPresentationRequest(BaseModel):
    title: str = Field(description="Presentation title")
    theme: str = Field(
        default='carbon_light',
        description=(
            "Visual theme. One of: "
            "'carbon_light' (white + IBM blue), "
            "'carbon_dark' (dark + electric blue), "
            "'consulting' (navy + cyan), "
            "'professional' (purple + magenta), "
            "'teal' (teal/green), "
            "'executive' (charcoal + orange)"
        )
    )
    slides: List[PremiumSlide] = Field(
        description="Ordered list of slides. Each must have a 'type' field."
    )


# ══════════════════════════════════════════════════════════════════════════════
# Tool
# ══════════════════════════════════════════════════════════════════════════════

@tool(
    name='create_premium_presentation',
    display_name='Create Premium IBM Presentation',
    description=(
        'Generate a polished IBM Carbon Design System presentation. '
        'Pixel-perfect backgrounds, header bands, accent stripes, and typography. '
        '6 themes: carbon_light, carbon_dark, consulting, professional, teal, executive. '
        'Slide types: title, content, section, two_column, stats, quote, thank_you. '
        'Returns a downloadable PPTX file. '
        'Use for premium-styled decks — distinct from the IBM template engine (create_presentation).'
    ),
)
def create_premium_presentation(request: PremiumPresentationRequest) -> bytes:
    """Generate a premium IBM Carbon-styled presentation and return PPTX bytes.

    Args:
        request: title, theme, and ordered list of slides.

    Returns:
        Binary PPTX bytes — WXO delivers this as a downloadable file.

    Raises:
        ValueError: If no slides provided or an unknown theme is given.
    """
    if not request.slides:
        raise ValueError('No slides provided.')

    if request.theme not in _THEMES:
        raise ValueError(
            f"Unknown theme '{request.theme}'. "
            f"Valid themes: {', '.join(_THEMES.keys())}"
        )

    gen = _IBMGen(request.theme)

    for i, s in enumerate(request.slides, 1):
        t = s.type

        if t == 'title':
            gen.add_title_slide(
                title    = s.title or request.title,
                subtitle = s.subtitle or '',
                date     = s.date or '',
                num      = i,
            )

        elif t == 'content':
            gen.add_content_slide(
                title   = s.title or '',
                content = _parse_list(s.content),
                eyebrow = s.eyebrow or '',
                layout  = s.layout or 'bullets',
                num     = i,
            )

        elif t == 'section':
            gen.add_section_divider(
                title   = s.title or '',
                eyebrow = s.eyebrow or '',
                num     = i,
            )

        elif t == 'two_column':
            gen.add_two_column_slide(
                title         = s.title or '',
                left_content  = _parse_list(s.left),
                right_content = _parse_list(s.right),
                left_label    = s.left_label or '',
                right_label   = s.right_label or '',
                num           = i,
            )

        elif t == 'stats':
            gen.add_stats_slide(
                title = s.title or '',
                stats = _parse_stats(s.stats),
                num   = i,
            )

        elif t == 'quote':
            gen.add_quote_slide(
                title       = s.title or '',
                quote       = s.quote or '',
                attribution = s.attribution or '',
                num         = i,
            )

        elif t == 'thank_you':
            gen.add_thank_you_slide(
                message = s.message or s.title or 'Thank You',
                contact = s.contact or '',
                num     = i,
            )

        # Unknown types are silently skipped — agent prompt already validates

    with tempfile.TemporaryDirectory() as tmp:
        safe = re.sub(r'[^\w\s\-]', '', request.title)
        safe = re.sub(r'\s+', '_', safe.strip())[:80] or 'presentation'
        out  = os.path.join(tmp, f'{safe}.pptx')
        gen.save(out)
        with open(out, 'rb') as fh:
            return fh.read()

# Made with Bob
