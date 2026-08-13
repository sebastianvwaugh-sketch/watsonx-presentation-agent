"""
IBM Presentation Generator — Premium Edition
Carbon Design System with enhanced visual polish.

Themes : carbon_light | carbon_dark | consulting | professional
Slides : title | content | section | two_column | stats | quote | thank_you
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Optional
import os

# ── IBM Carbon palette ────────────────────────────────────────────────────────
C = {
    # Blues
    'ibm_blue_60':    '0F62FE',
    'ibm_blue_70':    '0043CE',
    'ibm_blue_80':    '002D9C',
    'ibm_blue_40':    '4589FF',
    'ibm_blue_20':    'D0E2FF',
    'ibm_blue_10':    'EDF5FF',
    # Grays
    'gray_100':       '161616',
    'gray_90':        '262626',
    'gray_80':        '393939',
    'gray_70':        '525252',
    'gray_60':        '6F6F6F',
    'gray_40':        'A8A8A8',
    'gray_30':        'C6C6C6',
    'gray_20':        'E0E0E0',
    'gray_10':        'F4F4F4',
    # Accent colours
    'cyan_50':        '1192E8',
    'cyan_70':        '00539A',
    'cyan_20':        'BAE6FF',
    'teal_50':        '009D9A',
    'teal_70':        '004144',
    'teal_20':        '9EF0F0',
    'purple_50':      '8A3FFC',
    'purple_70':      '491D8B',
    'purple_20':      'E8DAFF',
    'magenta_50':     'EE5396',
    'magenta_70':     '9F1853',
    'green_50':       '24A148',
    'green_70':       '044317',
    # Neutrals
    'white':          'FFFFFF',
    'black':          '000000',
    # Warm
    'warm_gray_10':   'F7F3F2',
    'warm_gray_80':   '4D3D35',
    'red_60':         'DA1E28',
    'red_80':         '750E13',
    'orange_40':      'FF832B',
}

FONT = 'IBM Plex Sans'

_PROJECT_ROOT    = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DELIVERABLES_DIR = os.path.join(_PROJECT_ROOT, 'deliverables')


def _rgb(h: str) -> RGBColor:
    return RGBColor.from_string(h.lstrip('#'))


# ── Theme definitions ─────────────────────────────────────────────────────────
# Each theme has: bg, panel, panel2 (darker panel), accent, title_fg, body_fg,
#                 rule, footer_fg, tag_bg, tag_fg, bullet_color, stat_color

THEMES = {
    # Clean white with IBM Blue — classic consulting look
    'carbon_light': dict(
        bg=C['white'],          panel=C['ibm_blue_60'],    panel2=C['ibm_blue_70'],
        accent=C['ibm_blue_60'],title_fg=C['gray_100'],    body_fg=C['gray_70'],
        rule=C['gray_20'],      footer_fg=C['gray_60'],    tag_bg=C['ibm_blue_10'],
        tag_fg=C['ibm_blue_60'],bullet=C['ibm_blue_60'],   stat=C['ibm_blue_60'],
        cover_bg=C['ibm_blue_60'], cover_stripe=C['ibm_blue_80'],
    ),
    # Dark charcoal with electric blue — modern tech look
    'carbon_dark': dict(
        bg=C['gray_100'],       panel=C['gray_90'],         panel2=C['gray_80'],
        accent=C['ibm_blue_40'],title_fg=C['white'],         body_fg=C['gray_30'],
        rule=C['gray_80'],      footer_fg=C['gray_60'],     tag_bg=C['gray_80'],
        tag_fg=C['ibm_blue_40'],bullet=C['ibm_blue_40'],    stat=C['ibm_blue_40'],
        cover_bg=C['gray_90'],  cover_stripe=C['ibm_blue_60'],
    ),
    # Deep blue-navy with cyan — premium consulting
    'consulting': dict(
        bg=C['white'],          panel=C['ibm_blue_80'],     panel2=C['ibm_blue_70'],
        accent=C['cyan_50'],    title_fg=C['gray_100'],     body_fg=C['gray_70'],
        rule=C['gray_20'],      footer_fg=C['gray_60'],     tag_bg=C['cyan_20'],
        tag_fg=C['cyan_70'],    bullet=C['ibm_blue_80'],    stat=C['ibm_blue_80'],
        cover_bg=C['ibm_blue_80'], cover_stripe=C['cyan_50'],
    ),
    # Purple gradient feel — innovation / AI
    'professional': dict(
        bg=C['gray_10'],        panel=C['purple_50'],       panel2=C['purple_70'],
        accent=C['magenta_50'], title_fg=C['gray_100'],     body_fg=C['gray_70'],
        rule=C['gray_20'],      footer_fg=C['gray_60'],     tag_bg=C['purple_20'],
        tag_fg=C['purple_70'],  bullet=C['purple_50'],      stat=C['purple_50'],
        cover_bg=C['purple_70'], cover_stripe=C['magenta_50'],
    ),
    # Teal / green — sustainability / growth
    'teal': dict(
        bg=C['white'],          panel=C['teal_50'],         panel2=C['teal_70'],
        accent=C['teal_50'],    title_fg=C['gray_100'],     body_fg=C['gray_70'],
        rule=C['gray_20'],      footer_fg=C['gray_60'],     tag_bg=C['teal_20'],
        tag_fg=C['teal_70'],    bullet=C['teal_50'],        stat=C['teal_50'],
        cover_bg=C['teal_70'],  cover_stripe=C['teal_50'],
    ),
    # Warm charcoal — executive briefing
    'executive': dict(
        bg=C['warm_gray_10'],   panel=C['gray_100'],        panel2=C['gray_90'],
        accent=C['orange_40'],  title_fg=C['gray_100'],     body_fg=C['gray_70'],
        rule=C['gray_30'],      footer_fg=C['gray_60'],     tag_bg=C['gray_20'],
        tag_fg=C['gray_80'],    bullet=C['gray_100'],       stat=C['orange_40'],
        cover_bg=C['gray_100'], cover_stripe=C['orange_40'],
    ),
}


class IBMPresentationGenerator:
    W = Inches(13.33)
    H = Inches(7.5)

    def __init__(self, theme: str = 'carbon_light'):
        self.prs = Presentation()
        self.prs.slide_width  = self.W
        self.prs.slide_height = self.H
        self.tc = THEMES.get(theme, THEMES['carbon_light'])

    def _blank_slide(self):
        """Add a truly blank slide — removes all layout placeholders (slide number box etc.)."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # Remove every placeholder inherited from the layout (e.g. slide number box)
        for ph in slide.placeholders:
            sp = ph._element
            sp.getparent().remove(sp)
        return slide

    # ── Primitives ────────────────────────────────────────────────────────────

    def _rect(self, slide, x, y, w, h, fill, alpha=None):
        s = slide.shapes.add_shape(1, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
        s.line.fill.background()
        return s

    def _tb(self, slide, x, y, w, h, text, pt, bold=False,
            color='161616', align=PP_ALIGN.LEFT, wrap=True, italic=False, spacing=0):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tb.fill.background()          # explicit transparent fill — no black box
        tb.line.fill.background()     # no border line
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.name   = FONT
        p.font.size   = Pt(pt)
        p.font.bold   = bold
        p.font.italic = italic
        p.font.color.rgb = _rgb(color)
        if spacing:
            p.space_after = Pt(spacing)
        return tb

    def _footer(self, slide, num: int = None, light=False):
        tc = self.tc
        fg = C['white'] if light else tc['footer_fg']
        self._rect(slide, Inches(0.4), Inches(6.88), Inches(12.53), Inches(0.02),
                   C['white'] if light else tc['rule'])
        self._tb(slide, Inches(12.1), Inches(6.9), Inches(1.1), Inches(0.4),
                 'IBM', 16, bold=True, color=fg, align=PP_ALIGN.RIGHT)
        if num:
            self._tb(slide, Inches(0.4), Inches(6.9), Inches(0.8), Inches(0.4),
                     str(num), 10, color=fg)

    def _logo(self, slide, logo_path):
        if logo_path and os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path,
                Inches(11.6), Inches(6.65), height=Inches(0.55))

    # ── Cover slide ───────────────────────────────────────────────────────────

    def add_title_slide(self, title, subtitle='', logo_path=None, num=1):
        """
        Cover: full-bleed panel, left dark stripe, large left-aligned title,
        white rule divider, subtitle, IBM mark bottom-right.
        """
        slide = self._blank_slide()
        tc = self.tc

        # Full bleed
        self._rect(slide, 0, 0, self.W, self.H, tc['cover_bg'])

        # Left accent stripe (10% width)
        self._rect(slide, 0, 0, Inches(1.33), self.H, tc['cover_stripe'])

        # Subtle bottom bar
        self._rect(slide, 0, Inches(6.7), self.W, Inches(0.8), tc['panel2'])

        # White divider rule
        self._rect(slide, Inches(1.6), Inches(4.8), Inches(9.5), Inches(0.055), C['white'])

        # Title — large, bold, white
        self._tb(slide, Inches(1.6), Inches(1.8), Inches(10.5), Inches(2.7),
                 title, 56, bold=True, color=C['white'],
                 align=PP_ALIGN.LEFT, wrap=True)

        # Subtitle
        if subtitle:
            self._tb(slide, Inches(1.6), Inches(5.05), Inches(10.0), Inches(0.9),
                     subtitle, 22, color=C['white'], align=PP_ALIGN.LEFT)

        # IBM wordmark
        self._tb(slide, Inches(11.5), Inches(6.85), Inches(1.6), Inches(0.5),
                 'IBM', 20, bold=True, color=C['white'], align=PP_ALIGN.RIGHT)

        self._logo(slide, logo_path)
        return slide

    # ── Content slide ─────────────────────────────────────────────────────────

    def add_content_slide(self, title, content, layout='bullets', logo_path=None, num=2):
        """
        Content: bold header band, thin accent underline, generous bullet spacing,
        left accent stripe on header, slide number in footer.
        """
        slide = self._blank_slide()
        tc = self.tc

        # Slide background
        self._rect(slide, 0, 0, self.W, self.H, tc['bg'])

        # Header band
        self._rect(slide, 0, 0, self.W, Inches(1.45), tc['panel'])

        # Left accent notch in header
        self._rect(slide, 0, 0, Inches(0.22), Inches(1.45), tc['panel2'])

        # Accent underline
        self._rect(slide, 0, Inches(1.45), self.W, Inches(0.07), tc['accent'])

        # Title in header
        self._tb(slide, Inches(0.45), Inches(0.2), Inches(11.9), Inches(1.05),
                 title, 34, bold=True, color=C['white'], align=PP_ALIGN.LEFT)

        # Bullet body
        body_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.75), Inches(12.1), Inches(4.9))
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if layout == 'numbered':
                p.text = f"{i + 1}.  {item}"
                p.font.name  = FONT
                p.font.size  = Pt(20)
                p.font.color.rgb = _rgb(tc['bullet'])
            else:
                # Two runs: coloured dash + body text
                run_bullet = p.add_run()
                run_bullet.text = '\u2014  '
                run_bullet.font.name  = FONT
                run_bullet.font.size  = Pt(20)
                run_bullet.font.bold  = True
                run_bullet.font.color.rgb = _rgb(tc['bullet'])
                run_body = p.add_run()
                run_body.text = item
                run_body.font.name  = FONT
                run_body.font.size  = Pt(20)
                run_body.font.color.rgb = _rgb(tc['body_fg'])

            p.space_before = Pt(4)
            p.space_after  = Pt(10)

        self._footer(slide, num)
        self._logo(slide, logo_path)
        return slide

    # ── Section divider ───────────────────────────────────────────────────────

    def add_section_divider(self, title, eyebrow='', logo_path=None, num=2):
        """
        Section break: full panel, large left-aligned title, left stripe,
        subtle eyebrow label, decorative white accent block top-right.
        """
        slide = self._blank_slide()
        tc = self.tc

        self._rect(slide, 0, 0, self.W, self.H, tc['cover_bg'])

        # Left dark stripe
        self._rect(slide, 0, 0, Inches(1.33), self.H, tc['cover_stripe'])

        # Bottom bar
        self._rect(slide, 0, Inches(6.5), self.W, Inches(1.0), tc['panel2'])

        # Eyebrow
        label = eyebrow if eyebrow else 'SECTION'
        self._tb(slide, Inches(1.65), Inches(1.4), Inches(8.0), Inches(0.5),
                 label.upper(), 12, color=C['white'], italic=True)

        # Short accent rule below eyebrow
        self._rect(slide, Inches(1.65), Inches(1.95), Inches(1.2), Inches(0.055),
                   tc['cover_stripe'])

        # Main title
        self._tb(slide, Inches(1.65), Inches(2.1), Inches(9.2), Inches(3.0),
                 title, 52, bold=True, color=C['white'],
                 align=PP_ALIGN.LEFT, wrap=True)

        self._footer(slide, light=True)
        self._logo(slide, logo_path)
        return slide

    # ── Two-column ────────────────────────────────────────────────────────────

    def add_two_column_slide(self, title, left_content, right_content,
                             left_label='', right_label='',
                             logo_path=None, num=2):
        """
        Two-column: header band, labelled columns, vertical divider,
        per-column coloured sub-headers if labels provided.
        """
        slide = self._blank_slide()
        tc = self.tc

        self._rect(slide, 0, 0, self.W, self.H, tc['bg'])
        self._rect(slide, 0, 0, self.W, Inches(1.45), tc['panel'])
        self._rect(slide, 0, 0, Inches(0.22), Inches(1.45), tc['panel2'])
        self._rect(slide, 0, Inches(1.45), self.W, Inches(0.07), tc['accent'])

        self._tb(slide, Inches(0.45), Inches(0.2), Inches(11.9), Inches(1.05),
                 title, 34, bold=True, color=C['white'])

        # Vertical divider
        self._rect(slide, Inches(6.61), Inches(1.65), Inches(0.04), Inches(5.0), tc['rule'])

        col_top = Inches(1.65)

        for items, label, x in [
            (left_content,  left_label,  Inches(0.55)),
            (right_content, right_label, Inches(6.85)),
        ]:
            cur_top = col_top
            if label:
                # Coloured label pill
                self._rect(slide, x, cur_top, Inches(5.7), Inches(0.42), tc['tag_bg'])
                self._tb(slide, x + Inches(0.12), cur_top + Inches(0.04),
                         Inches(5.4), Inches(0.34),
                         label, 14, bold=True, color=tc['tag_fg'])
                cur_top += Inches(0.55)

            cb = slide.shapes.add_textbox(x, cur_top, Inches(5.7), Inches(4.5))
            cf = cb.text_frame
            cf.word_wrap = True
            for i, item in enumerate(items[:6]):
                p = cf.paragraphs[0] if i == 0 else cf.add_paragraph()
                r1 = p.add_run()
                r1.text = '\u2014  '
                r1.font.name = FONT
                r1.font.size = Pt(18)
                r1.font.bold = True
                r1.font.color.rgb = _rgb(tc['bullet'])
                r2 = p.add_run()
                r2.text = item
                r2.font.name = FONT
                r2.font.size = Pt(18)
                r2.font.color.rgb = _rgb(tc['body_fg'])
                p.space_after = Pt(9)

        self._footer(slide, num)
        self._logo(slide, logo_path)
        return slide

    # ── Stats slide ───────────────────────────────────────────────────────────

    def add_stats_slide(self, title, stats, logo_path=None, num=2):
        """
        Stats: 2–4 large headline numbers with labels, separated by vertical rules.
        stats = [{"number": "85%", "label": "Workloads migrated"}, ...]
        """
        slide = self._blank_slide()
        tc = self.tc

        self._rect(slide, 0, 0, self.W, self.H, tc['bg'])
        self._rect(slide, 0, 0, self.W, Inches(1.45), tc['panel'])
        self._rect(slide, 0, 0, Inches(0.22), Inches(1.45), tc['panel2'])
        self._rect(slide, 0, Inches(1.45), self.W, Inches(0.07), tc['accent'])
        self._tb(slide, Inches(0.45), Inches(0.2), Inches(11.9), Inches(1.05),
                 title, 34, bold=True, color=C['white'])

        n = min(len(stats), 4)
        col_w = 12.53 / n
        for i, stat in enumerate(stats[:4]):
            x = Inches(0.4 + i * col_w)
            w = Inches(col_w)

            # Vertical divider (not before first)
            if i > 0:
                self._rect(slide, x - Inches(0.02), Inches(1.8),
                           Inches(0.04), Inches(4.6), tc['rule'])

            # Big number
            self._tb(slide, x + Inches(0.15), Inches(2.1), w - Inches(0.3), Inches(2.2),
                     stat.get('number', '—'), 72, bold=True,
                     color=tc['stat'], align=PP_ALIGN.LEFT)

            # Accent underline below number
            self._rect(slide, x + Inches(0.15), Inches(3.9),
                       Inches(col_w * 0.35), Inches(0.055), tc['accent'])

            # Label
            self._tb(slide, x + Inches(0.15), Inches(4.1), w - Inches(0.3), Inches(2.3),
                     stat.get('label', ''), 18,
                     color=tc['body_fg'], wrap=True)

            # Source / footnote
            if stat.get('source'):
                self._tb(slide, x + Inches(0.15), Inches(6.2), w - Inches(0.3), Inches(0.5),
                         stat['source'], 11, italic=True, color=tc['footer_fg'])

        self._footer(slide, num)
        self._logo(slide, logo_path)
        return slide

    # ── Quote slide ───────────────────────────────────────────────────────────

    def add_quote_slide(self, title, quote, attribution='', logo_path=None, num=2):
        """
        Quote: large left-hand opening quote mark, quote text, attribution line.
        """
        slide = self._blank_slide()
        tc = self.tc

        self._rect(slide, 0, 0, self.W, self.H, tc['bg'])

        # Left accent stripe
        self._rect(slide, 0, 0, Inches(0.45), self.H, tc['panel'])

        # Slide title — top of slide, clear of quote mark
        self._tb(slide, Inches(0.7), Inches(0.15), Inches(11.5), Inches(0.6),
                 title, 22, bold=True, color=tc['accent'])

        # Short accent rule below title
        self._rect(slide, Inches(0.7), Inches(0.78), Inches(1.5), Inches(0.055), tc['accent'])

        # Giant decorative quote mark — below title
        self._tb(slide, Inches(0.7), Inches(0.95), Inches(2.0), Inches(2.0),
                 '\u201C', 120, bold=True, color=tc['accent'])

        # Quote text
        self._tb(slide, Inches(0.8), Inches(2.1), Inches(11.3), Inches(3.6),
                 quote, 26, bold=False, color=tc['title_fg'],
                 align=PP_ALIGN.LEFT, wrap=True, italic=True)

        # Attribution
        if attribution:
            self._rect(slide, Inches(0.7), Inches(5.95), Inches(0.4), Inches(0.55), tc['accent'])
            self._tb(slide, Inches(1.25), Inches(5.95), Inches(11.0), Inches(0.55),
                     attribution, 16, bold=True, color=tc['title_fg'])

        self._footer(slide, num)
        self._logo(slide, logo_path)
        return slide

    # ── Thank you ─────────────────────────────────────────────────────────────

    def add_thank_you_slide(self, message='Thank You', contact='', logo_path=None, num=2):
        """
        Closing: full-bleed panel, large left-aligned message, rule,
        contact detail, left stripe.
        """
        slide = self._blank_slide()
        tc = self.tc

        self._rect(slide, 0, 0, self.W, self.H, tc['cover_bg'])
        self._rect(slide, 0, 0, Inches(1.33), self.H, tc['cover_stripe'])
        self._rect(slide, 0, Inches(6.5), self.W, Inches(1.0), tc['panel2'])
        self._rect(slide, Inches(1.6), Inches(5.0), Inches(9.5), Inches(0.055), C['white'])

        self._tb(slide, Inches(1.6), Inches(2.2), Inches(10.5), Inches(2.5),
                 message, 60, bold=True, color=C['white'], wrap=True)

        if contact:
            self._tb(slide, Inches(1.6), Inches(5.2), Inches(10.0), Inches(0.9),
                     contact, 20, color=C['white'])

        self._tb(slide, Inches(11.5), Inches(6.85), Inches(1.6), Inches(0.5),
                 'IBM', 20, bold=True, color=C['white'], align=PP_ALIGN.RIGHT)

        self._logo(slide, logo_path)
        return slide

    def save(self, output_path: str) -> str:
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        self.prs.save(output_path)
        return output_path


# ── Public API ────────────────────────────────────────────────────────────────

def create_ibm_presentation(
    title: str,
    slides: List[Dict],
    theme: str = 'carbon_light',
    output_path: Optional[str] = None,
    logo_path: Optional[str] = None,
) -> str:
    """
    Create an IBM-branded presentation and save it.

    Themes: carbon_light | carbon_dark | consulting | professional | teal | executive

    Slide types:
        {"type": "title",      "title": "...", "subtitle": "..."}
        {"type": "content",    "title": "...", "content": [...], "layout": "bullets|numbered"}
        {"type": "section",    "title": "...", "eyebrow": "..."}
        {"type": "two_column", "title": "...", "left": [...], "right": [...],
                               "left_label": "...", "right_label": "..."}
        {"type": "stats",      "title": "...",
                               "stats": [{"number":"85%","label":"...","source":"..."}]}
        {"type": "quote",      "title": "...", "quote": "...", "attribution": "..."}
        {"type": "thank_you",  "message": "...", "contact": "..."}
    """
    if logo_path is None:
        sources_dir = os.path.join(_PROJECT_ROOT, 'sources')
        if os.path.isdir(sources_dir):
            for fname in sorted(os.listdir(sources_dir)):
                if 'logo' in fname.lower() and fname.lower().endswith(('.png', '.jpg', '.jpeg')):
                    logo_path = os.path.join(sources_dir, fname)
                    break

    if output_path is None:
        safe = title.replace(' ', '_').replace('/', '_')
        output_path = os.path.join(DELIVERABLES_DIR, f"{safe}.pptx")

    gen = IBMPresentationGenerator(theme)

    for i, sd in enumerate(slides, 1):
        t = sd.get('type', 'content')
        if t == 'title':
            gen.add_title_slide(sd.get('title', title), sd.get('subtitle', ''), logo_path, i)
        elif t == 'content':
            gen.add_content_slide(sd['title'], sd.get('content', []), sd.get('layout', 'bullets'), logo_path, i)
        elif t == 'section':
            gen.add_section_divider(sd['title'], sd.get('eyebrow', ''), logo_path, i)
        elif t == 'two_column':
            gen.add_two_column_slide(sd['title'], sd.get('left', []), sd.get('right', []),
                                     sd.get('left_label', ''), sd.get('right_label', ''), logo_path, i)
        elif t == 'stats':
            gen.add_stats_slide(sd['title'], sd.get('stats', []), logo_path, i)
        elif t == 'quote':
            gen.add_quote_slide(sd['title'], sd.get('quote', ''), sd.get('attribution', ''), logo_path, i)
        elif t == 'thank_you':
            gen.add_thank_you_slide(sd.get('message', 'Thank You'), sd.get('contact', ''), logo_path, i)

    return gen.save(output_path)


# ── Smoke test — generates all 6 themes ──────────────────────────────────────

if __name__ == '__main__':
    demo_slides = [
        {'type': 'title', 'title': 'AI Transformation at Scale', 'subtitle': 'IBM Consulting — Strategic Roadmap 2025'},
        {'type': 'stats', 'title': 'The Case for Action', 'stats': [
            {'number': '72%', 'label': 'of enterprises say AI is a top-3 priority in 2025', 'source': 'IBM IBV 2024'},
            {'number': '3.4x', 'label': 'ROI delivered by organisations with mature AI programmes', 'source': 'IBM IBV 2024'},
            {'number': '$18B', 'label': 'estimated productivity gain from AI in financial services by 2026', 'source': 'Gartner 2024'},
        ]},
        {'type': 'content', 'title': 'Key Strategic Priorities', 'content': [
            'Embed AI into core business processes — not just point solutions',
            'Build internal AI capability alongside vendor partnerships',
            'Establish data governance before scaling AI deployments',
            'Measure outcomes in business terms: cost, speed, quality, risk',
        ]},
        {'type': 'section', 'title': 'Our Recommended Approach', 'eyebrow': 'Part 2'},
        {'type': 'two_column', 'title': 'Current State vs Target State',
         'left_label': 'Where You Are Today', 'right_label': 'Where You Need to Be',
         'left':  ['Siloed AI experiments across BUs', 'No unified data platform', 'Manual governance processes', 'Limited AI talent pipeline'],
         'right': ['Enterprise-wide AI programme office', 'Governed data mesh architecture', 'Automated compliance & audit trails', 'AI Centre of Excellence operational']},
        {'type': 'quote', 'title': 'Executive Perspective',
         'quote': 'The organisations that will win are not those with the best AI models — they are those who embed AI into how they make decisions every single day.',
         'attribution': 'Arvind Krishna, Chairman & CEO, IBM'},
        {'type': 'content', 'title': 'Immediate Next Steps', 'content': [
            'Week 1–2: AI readiness assessment across top 5 business units',
            'Week 3–4: Identify 3 high-value pilot use cases with measurable KPIs',
            'Month 2: Stand up data governance working group',
            'Month 3: Launch first pilot with dedicated IBM Consulting team',
            'Month 6: Review outcomes and define scale-up plan',
        ]},
        {'type': 'thank_you', 'message': 'Let\'s Build This Together', 'contact': 'sebastian.waugh@ibm.com  |  +44 7700 900000'},
    ]

    for theme in ['carbon_light', 'carbon_dark', 'consulting', 'professional', 'teal', 'executive']:
        path = create_ibm_presentation(
            'AI Transformation at Scale',
            demo_slides,
            theme=theme,
            output_path=f'deliverables/premium_{theme}.pptx'
        )
        size = round(os.path.getsize(path) / 1024, 1)
        print(f'  {theme:<18}  {size} KB  ->  {path}')
