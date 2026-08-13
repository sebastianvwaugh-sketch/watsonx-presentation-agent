#!/usr/bin/env python3
"""pptx_workflow.py (request-first)

This script generates a PPTX deterministically from a small request.json.

Fixes included
- Robust bullet parsing for request strings (prevents slide 21 body being cleared).
- Supports indexed value shapes without trailing underscore (e.g., stat_01) via macro arrays.
- Indexed-family deletion removes both lead_XX_* and exact lead_XX value shapes.
- Optional request.meta.expected_slide_count warns if deck size differs.

Diagnostics
- Default: print issues to stdout; illegal-target payload to stderr and exit 2.
- --report writes JSON to a file (optional).
"""

from __future__ import annotations

import argparse
import copy
import json
import os
import re
import sys
import zipfile
from typing import Any, Dict, List, Optional, Tuple


def parse_markdown_table(table_text: str):
    """Parse a GitHub-flavored markdown pipe table to {'headers': [...], 'rows': [[...], ...]}.

    Supports:
      | H1 | H2 |
      | --- | --- |
      | a | b |

    Returns empty headers/rows if parsing fails.
    """
    if not isinstance(table_text, str) or not table_text.strip():
        return {'headers': [], 'rows': []}
    lines = [ln.strip() for ln in table_text.strip().splitlines() if ln.strip()]
    pipe_lines = [ln for ln in lines if '|' in ln]
    if len(pipe_lines) < 2:
        return {'headers': [], 'rows': []}

    def split_row(ln: str):
        ln = ln.strip()
        if ln.startswith('|'):
            ln = ln[1:]
        if ln.endswith('|'):
            ln = ln[:-1]
        return [c.strip() for c in ln.split('|')]

    header = split_row(pipe_lines[0])
    sep = pipe_lines[1]
    if re.fullmatch(r'\|?\s*[:\-\s\|]+\s*\|?', sep):
        body_lines = pipe_lines[2:]
    else:
        body_lines = pipe_lines[1:]
    rows = [split_row(ln) for ln in body_lines]
    n = max([len(header)] + [len(r) for r in rows] + [0])
    header = (header + [''] * n)[:n]
    rows = [(r + [''] * n)[:n] for r in rows]
    return {'headers': header, 'rows': rows}

import lxml.etree as ET
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


def iter_shapes(slide):
    for shp in slide.shapes:
        yield shp
        try:
            if getattr(shp, 'shape_type', None) == 6:  # group
                for sub in shp.shapes:
                    yield sub
        except Exception:
            pass


def get_shape_by_name(slide, name: str):
    for shp in iter_shapes(slide):
        if (getattr(shp, 'name', None) or '') == name:
            return shp
    raise KeyError(f"Shape named '{name}' not found")


def shape_kind(shape) -> str:
    if getattr(shape, 'has_table', False):
        return 'table'
    if getattr(shape, 'has_text_frame', False):
        return 'text'
    try:
        if getattr(shape, 'shape_type', None) == 13:
            return 'picture'
    except Exception:
        pass
    try:
        blips = shape._element.xpath('.//*[local-name()="blip" or local-name()="svgBlip"]')
        if blips:
            return 'picture'
    except Exception:
        pass
    return 'other'


def get_slide_notes_text(slide) -> str:
    try:
        ns = slide.notes_slide
        if ns and ns.notes_text_frame:
            return ns.notes_text_frame.text or ''
    except Exception:
        return ''
    return ''


def extract_notes_markdown(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    out: List[str] = [f"# Speaker notes — {os.path.basename(pptx_path)}", ""]
    for i, slide in enumerate(prs.slides, start=1):
        notes = (get_slide_notes_text(slide) or '').strip()
        out.append(f"## Slide {i:02d}")
        out.append('')
        out.append(notes if notes else '(no speaker notes)')
        out.append('')
    return '\n'.join(out)


def _compile_shape_line_regex(shape_name: str) -> re.Pattern:
    sn = re.escape(shape_name)
    return re.compile(rf"(^|[\s\-\*•])({sn})(?=\s|:|$)")


def _line_mentions_shape(line: str, shape_name: str) -> bool:
    if _compile_shape_line_regex(shape_name).search(line):
        return True
    tokens = re.findall(r"[A-Za-z0-9_]+", line)
    sname = shape_name.strip()
    sname_l = sname.lower()
    for tok in tokens:
        tl = tok.lower()
        if '0x' in tl:
            pat = re.escape(tok)
            pat = re.sub(r"0x", lambda _m: r"\d{2}", pat, flags=re.I)
            if re.fullmatch(pat, sname):
                return True
        if tl == 'table':
            if sname_l == 'table' or re.fullmatch(r'table_\d{2}', sname_l):
                return True
    return False


def find_notes_line_for_shape(notes_lines: List[str], shape_name: str) -> Optional[str]:
    matches = [ln.strip() for ln in notes_lines if _line_mentions_shape(ln, shape_name)]
    if not matches:
        return None
    def score(line: str) -> float:
        l = line.lower()
        s = 0.0
        for kw, pts in [
            ('style guardrails', 5), ('mixed color', 5), ('mixed formatting', 5),
            ('do not', 4), ('avoid', 3), ('safe edit', 3), ('formatting-sensitive', 3), ('mandatory', 2)
        ]:
            if kw in l:
                s += pts
        s += 1.0 / max(1, len(line))
        return s
    matches.sort(key=score, reverse=True)
    return matches[0]


def paragraph_plain_text(p) -> str:
    return ''.join(run.text or '' for run in p.runs)


def has_bullet(paragraph) -> bool:
    ppr = paragraph._p.find(qn('a:pPr'))
    if ppr is None:
        return False
    for child in list(ppr):
        local = child.tag.split('}')[-1]
        if local in ('buChar', 'buAutoNum', 'buBlip'):
            return True
    return False


def extract_default_for_text_shape(shape) -> Tuple[str, Any]:
    tf = shape.text_frame
    paras = tf.paragraphs
    if any(has_bullet(p) for p in paras):
        bullets = []
        for p in paras:
            if has_bullet(p):
                t = paragraph_plain_text(p).strip()
                if t:
                    bullets.append(t)
        return 'set_bullets', bullets
    texts = [paragraph_plain_text(p).strip() for p in paras]
    trimmed = [t for t in texts if t != '']
    if len(trimmed) > 1:
        return 'set_paragraphs', trimmed
    return 'set_text', '\n'.join([t for t in texts if t])


def extract_default_for_table_shape(shape) -> Dict[str, Any]:
    tbl = shape.table
    rows = len(tbl.rows)
    cols = len(tbl.columns)
    cells: List[List[str]] = []
    for r in range(rows):
        row: List[str] = []
        for c in range(cols):
            tf = tbl.cell(r, c).text_frame
            row.append((tf.text or '') if tf is not None else '')
        cells.append(row)
    return {'rows': rows, 'cols': cols, 'cells': cells}


def allowed_shapes_for_slide(slide, notes_text: str) -> Dict[str, Dict[str, Any]]:
    lines = [ln.rstrip() for ln in (notes_text or '').splitlines() if ln.strip()]
    out: Dict[str, Dict[str, Any]] = {}
    for shp in iter_shapes(slide):
        nm = (getattr(shp, 'name', '') or '').strip()
        if not nm:
            continue
        hint = find_notes_line_for_shape(lines, nm)
        if not hint:
            continue
        k = shape_kind(shp)
        if k == 'text':
            op, default = extract_default_for_text_shape(shp)
        elif k == 'table':
            op, default = 'set_table', extract_default_for_table_shape(shp)
        elif k == 'picture':
            op, default = 'swap_picture', None
        else:
            continue
        out[nm] = {'kind': k, 'op': op, 'default': default, 'notes_hint': hint}
    return out


def parse_indexed_family(shape_name: str) -> Optional[Tuple[str, str]]:
    m = re.match(r'^([a-zA-Z]+_)(\d{2})_', shape_name)
    if not m:
        return None
    return m.group(1), m.group(2)


def parse_indexed_value_shape(shape_name: str) -> Optional[Tuple[str, str]]:
    m = re.match(r'^([a-zA-Z]+_)(\d{2})$', shape_name)
    if not m:
        return None
    return m.group(1), m.group(2)


def parse_indexed_any(shape_name: str) -> Optional[Tuple[str, str]]:
    return parse_indexed_family(shape_name) or parse_indexed_value_shape(shape_name)


def plural_macro_key(lead_: str) -> str:
    base = lead_[:-1] if lead_.endswith('_') else lead_
    return base + 's'


def is_empty_value(op: str, value: Any) -> bool:
    if value is None:
        return True
    if op == 'set_text':
        return (not isinstance(value, str)) or (value.strip() == '')
    if op in ('set_bullets', 'set_paragraphs'):
        if not isinstance(value, list):
            return True
        if len(value) == 0:
            return True
        return all((str(x).strip() == '') for x in value)
    if op == 'set_table':
        if not isinstance(value, dict):
            return True
        if value.get('headers') or value.get('rows'):
            return False
        # Legacy format: {'cells': [[...], ...]}
        cells = value.get('cells')
        if isinstance(cells, list) and any(any(str(x).strip() for x in row) for row in cells if isinstance(row, list)):
            return False
        return True
    if op == 'swap_picture':
        if not isinstance(value, dict):
            return True
        return (not value.get('asset_slide')) or (not (value.get('asset_shape') or '').strip())
    if isinstance(value, str):
        return value.strip() == ''
    if isinstance(value, list):
        return len(value) == 0 or all((str(x).strip() == '') for x in value)
    if isinstance(value, dict):
        return len(value) == 0
    return False


def compute_family_delete_rules(values_by_shape: Dict[str, Dict[str, Any]]) -> Tuple[List[str], List[str], List[Dict[str, Any]]]:
    by_group: Dict[Tuple[str, str], List[str]] = {}
    for nm in values_by_shape.keys():
        fam = parse_indexed_any(nm)
        if not fam:
            continue
        lead, idx = fam
        by_group.setdefault((lead, idx), []).append(nm)

    group_empty: Dict[Tuple[str, str], bool] = {}
    for (lead, idx), shapes in by_group.items():
        group_empty[(lead, idx)] = all(is_empty_value(values_by_shape[n]['op'], values_by_shape[n]['value']) for n in shapes)

    by_lead: Dict[str, List[str]] = {}
    for (lead, idx) in by_group.keys():
        by_lead.setdefault(lead, []).append(idx)

    delete_prefixes: List[str] = []
    delete_exact: List[str] = []
    warnings: List[Dict[str, Any]] = []

    for lead, idxs in by_lead.items():
        u = sorted(set(idxs))
        filled = [i for i in u if not group_empty.get((lead, i), True)]
        max_filled = max(filled) if filled else None
        for idx in u:
            if group_empty.get((lead, idx), True):
                delete_prefixes.append(f"{lead}{idx}_")
                delete_exact.append(f"{lead}{idx}")
                if max_filled and idx < max_filled:
                    warnings.append({'lead': lead, 'empty_index': idx, 'max_filled_index': max_filled})

    return sorted(set(delete_prefixes)), sorted(set(delete_exact)), warnings


# ---- Robust bullet parsing ----

def parse_bullets_from_string(text: str) -> List[str]:
    lines = [ln.rstrip() for ln in (text or '').splitlines()]
    bullets: List[str] = []
    cur: List[str] = []
    start_re = re.compile(r'^[\s\ufeff\u200b]*(?:[•\-*\u2013\u2014\u2212])\s*(.*)$')

    for ln in lines:
        if not ln.strip():
            continue
        m = start_re.match(ln)
        if m:
            if cur:
                bullets.append(' '.join(cur).strip())
                cur = []
            cur.append((m.group(1) or '').strip())
        else:
            if not cur:
                cur.append(ln.strip())
            else:
                cur.append(ln.strip())

    if cur:
        bullets.append(' '.join(cur).strip())

    bullets = [b for b in bullets if b]
    if not bullets and (text or '').strip():
        bullets = [ln.strip() for ln in (text or '').splitlines() if ln.strip()]
    return bullets


def coerce_value_for_op(op: str, provided: Any, shape_name: str) -> Any:
    if op == 'set_text':
        if provided is None:
            return ''
        if isinstance(provided, list):
            return '\n'.join(str(x) for x in provided)
        return str(provided)

    if op == 'set_table':
        if provided is None:
            return {'headers': [], 'rows': []}
        if isinstance(provided, dict):
            return provided
        if isinstance(provided, str):
            parsed = parse_markdown_table(provided)
            return parsed if (parsed.get('headers') or parsed.get('rows')) else {'headers': [], 'rows': []}
        if isinstance(provided, list):
            # list-of-lists treated as rows
            if all(isinstance(r, list) for r in provided):
                return {'headers': [], 'rows': [[str(x) for x in r] for r in provided]}
            return {'headers': [], 'rows': [[str(x)] for x in provided]}
        return {'headers': [], 'rows': [[str(provided)]]}

    if op in ('set_paragraphs', 'set_bullets'):
        if provided is None:
            return []
        if isinstance(provided, str):
            if op == 'set_bullets' and shape_name != 'agenda_list':
                parsed = parse_bullets_from_string(provided)
                return parsed if parsed else [ln.strip() for ln in provided.splitlines() if ln.strip()]
            return [ln.strip() for ln in provided.splitlines() if ln.strip()]
        if isinstance(provided, list):
            return [str(x).strip() for x in provided if str(x).strip()]
        return [str(provided)]

    return provided


def asset_ref_from_request(val: Any, asset_slide: Optional[int]) -> Dict[str, Any]:
    if isinstance(val, dict):
        if 'asset_slide' in val and 'asset_shape' in val:
            return {'asset_slide': val.get('asset_slide'), 'asset_shape': val.get('asset_shape', '')}
        if 'asset' in val:
            return {'asset_slide': asset_slide, 'asset_shape': val.get('asset', '')}
    if isinstance(val, str):
        return {'asset_slide': asset_slide, 'asset_shape': val}
    return {'asset_slide': asset_slide, 'asset_shape': ''}


def expand_request_to_plan(prs: Presentation, template_path: str, request: Dict[str, Any]) -> Tuple[Dict[str, Any], List[Dict[str, Any]], Dict[str, Any]]:
    deck_req = request.get('deck') or []
    content = request.get('content') or {}
    if not isinstance(deck_req, list):
        raise ValueError('request.deck must be a list')

    # deck size expectation warning
    meta = request.get('meta') or {}
    if isinstance(meta, dict) and isinstance(meta.get('expected_slide_count'), int):
        exp_n = int(meta.get('expected_slide_count'))
        if exp_n and exp_n != len(deck_req):
            print(f"[pptx_workflow] Warning: request.meta.expected_slide_count={exp_n} but request.deck has {len(deck_req)} entries.", file=sys.stderr)

    # Asset slides (best-effort via notes)
    asset_image_slide = None
    asset_pict_slide = None
    for i, slide in enumerate(prs.slides, start=1):
        notes = get_slide_notes_text(slide)
        if 'Asset slide: image_background' in notes:
            asset_image_slide = i
        if 'Asset slide: Pictogram' in notes or 'Asset slide: Pictogram option library' in notes:
            asset_pict_slide = i

    plan = {
        'template': {'file': os.path.basename(template_path)},
        'assets': {'image_library_slide': asset_image_slide, 'pictogram_library_slide': asset_pict_slide},
        'deck': [],
        'edits': {},
    }

    warnings: List[Dict[str, Any]] = []
    allowed_report: Dict[str, Any] = {}
    illegal: List[Dict[str, Any]] = []

    for item in deck_req:
        slide_no = int(item['slide_no'])
        sid = str(item['id'])
        plan['deck'].append({'id': sid, 'source_slide': slide_no})

        slide = prs.slides[slide_no - 1]
        notes = get_slide_notes_text(slide) or ''
        allowed = allowed_shapes_for_slide(slide, notes)
        allowed_names = sorted(allowed.keys())

        allowed_report[sid] = {
            'source_slide': slide_no,
            'allowed_shapes': [
                {'shape': nm, 'op': allowed[nm]['op'], 'kind': allowed[nm]['kind'], 'notes_hint': allowed[nm]['notes_hint']}
                for nm in allowed_names
            ],
        }

        slide_content = (content.get(sid, {}) or {})
        if not isinstance(slide_content, dict):
            raise ValueError(f"request.content['{sid}'] must be an object")

        macro_keys = {'items', 'steps', 'stats'}
        for k in slide_content.keys():
            if k in macro_keys or k == 'agenda_list':
                continue
            if k not in allowed:
                illegal.append({'id': sid, 'source_slide': slide_no, 'illegal_shape': k})

        items = slide_content.get('items') if isinstance(slide_content.get('items'), list) else None
        steps = slide_content.get('steps') if isinstance(slide_content.get('steps'), list) else None
        stats = slide_content.get('stats') if isinstance(slide_content.get('stats'), list) else None

        slide_edits: List[Dict[str, Any]] = []

        for nm in allowed_names:
            meta2 = allowed[nm]
            op = meta2['op']
            default = meta2['default']
            knd = meta2['kind']

            value: Any = None

            if nm in slide_content:
                provided = slide_content[nm]
                if op == 'swap_picture':
                    asset_slide = asset_pict_slide if 'pictogram' in nm else asset_image_slide
                    value = asset_ref_from_request(provided, asset_slide)
                elif op == 'set_table':
                    value = coerce_value_for_op(op, provided, nm)
                else:
                    value = coerce_value_for_op(op, provided, nm)

            elif nm == 'agenda_list' and 'agenda_list' in slide_content:
                value = coerce_value_for_op(op, slide_content['agenda_list'], nm)

            # fill lead_XX value shapes from macro array pluralised, e.g. stat_01 from stats[0].value
            elif re.fullmatch(r'[A-Za-z]+_\d{2}', nm):
                famv = parse_indexed_value_shape(nm)
                if famv:
                    lead_, idx = famv
                    arr = slide_content.get(plural_macro_key(lead_))
                    i = int(idx) - 1
                    if isinstance(arr, list) and 0 <= i < len(arr) and isinstance(arr[i], dict) and 'value' in arr[i]:
                        value = coerce_value_for_op(op, arr[i].get('value'), nm)

            else:
                fam = parse_indexed_family(nm)
                if fam:
                    lead, idx = fam
                    i = int(idx) - 1

                    if lead == 'item_' and items is not None and 0 <= i < len(items) and isinstance(items[i], dict):
                        it = items[i]
                        if nm.endswith('_title'):
                            value = coerce_value_for_op(op, it.get('title'), nm)
                        elif nm.endswith('_heading'):
                            value = coerce_value_for_op(op, it.get('heading'), nm)
                        elif nm.endswith('_body'):
                            provided = it.get('body_paragraphs') if 'body_paragraphs' in it else it.get('body')
                            value = coerce_value_for_op(op, provided, nm)
                        elif nm.endswith('_pictogram'):
                            value = asset_ref_from_request(it.get('pictogram'), asset_pict_slide)

                    if lead == 'step_' and steps is not None and 0 <= i < len(steps) and isinstance(steps[i], dict):
                        st = steps[i]
                        if nm.endswith('_body'):
                            provided = st.get('body_lines') if 'body_lines' in st else st.get('body')
                            value = coerce_value_for_op(op, provided, nm)
                        elif nm.endswith('_pictogram'):
                            value = asset_ref_from_request(st.get('pictogram'), asset_pict_slide)

                    if lead == 'stat_' and stats is not None and 0 <= i < len(stats) and isinstance(stats[i], dict):
                        stt = stats[i]
                        if nm.endswith('_eyebrow'):
                            value = coerce_value_for_op(op, stt.get('eyebrow'), nm)
                        elif nm.endswith('_body'):
                            value = coerce_value_for_op(op, stt.get('body'), nm)
                        elif nm.endswith('_pictogram'):
                            value = asset_ref_from_request(stt.get('pictogram'), asset_pict_slide)

            if value is None:
                if op == 'set_text':
                    value = ''
                elif op in ('set_bullets', 'set_paragraphs'):
                    value = []
                elif op == 'set_table':
                    value = {'headers': [], 'rows': []}
                elif op == 'swap_picture':
                    value = {'asset_slide': None, 'asset_shape': ''}

            slide_edits.append({'op': op, 'shape': nm, 'value': value, 'default': default, 'notes_hint': meta2.get('notes_hint',''), 'kind': knd})

        plan['edits'][sid] = slide_edits

    report = {'allowed': allowed_report}
    if illegal:
        report['illegal'] = illegal

    return plan, warnings, report


# -----------------------------
# In-place text editing
# -----------------------------

def clear_paragraph(p):
    for r in p.runs:
        r.text = ''


def set_paragraph_text(p, text: str):
    # agenda mixed runs: preserve number prefix run
    if p.runs and len(p.runs) >= 2:
        r0 = (p.runs[0].text or '')
        if re.match(r'^\d{2}\.\s*$', r0):
            m2 = re.match(r'^(\d{2}\.\s*)(.*)$', text)
            if m2:
                p.runs[0].text = m2.group(1)
                p.runs[1].text = m2.group(2)
            else:
                p.runs[1].text = text
            for r in p.runs[2:]:
                r.text = ''
            return

    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ''
    else:
        r = p.add_run(); r.text = text


def clone_paragraph_after(paragraph):
    new_p = copy.deepcopy(paragraph._p)
    paragraph._p.addnext(new_p)


def set_shape_text_single_paragraph(shape, text: str):
    tf = shape.text_frame
    if not tf.paragraphs:
        _ = tf.paragraphs
    set_paragraph_text(tf.paragraphs[0], text)
    for p in tf.paragraphs[1:]:
        clear_paragraph(p)


def set_multi_paragraph(shape, paragraphs_text: List[str]):
    tf = shape.text_frame
    if not tf.paragraphs:
        _ = tf.paragraphs
    while len(tf.paragraphs) < len(paragraphs_text):
        clone_paragraph_after(tf.paragraphs[-1])
    paras = tf.paragraphs
    for i, t in enumerate(paragraphs_text):
        set_paragraph_text(paras[i], t)
    for p in paras[len(paragraphs_text):]:
        clear_paragraph(p)


def set_bullets_preserve_style(shape, bullets: List[str]):
    tf = shape.text_frame
    bullet_paras = [p for p in tf.paragraphs if has_bullet(p)]
    if not bullet_paras:
        bullet_paras = tf.paragraphs

    while len(bullet_paras) < len(bullets):
        clone_paragraph_after(bullet_paras[-1])
        bullet_paras = [p for p in tf.paragraphs if has_bullet(p)] or tf.paragraphs

    for i, txt in enumerate(bullets):
        set_paragraph_text(bullet_paras[i], txt)

    for p in bullet_paras[len(bullets):]:
        clear_paragraph(p)

    # IMPORTANT: do not rely on Paragraph equality for membership checks.
    keep = {p._p for p in bullet_paras}
    for p in tf.paragraphs:
        if p._p in keep:
            continue
        clear_paragraph(p)

def set_table_safe(slide, table_shape_name: str, headers: List[str], rows: List[List[str]]):
    shp = get_shape_by_name(slide, table_shape_name)
    tbl = shp.table
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)

    for j in range(col_count):
        text = headers[j] if j < len(headers) else ''
        tf = tbl.cell(0, j).text_frame
        set_paragraph_text(tf.paragraphs[0], text)
        for p in tf.paragraphs[1:]:
            clear_paragraph(p)

    body_capacity = row_count - 1
    use_rows = rows[:body_capacity]

    for i in range(1, row_count):
        r_idx = i - 1
        for j in range(col_count):
            val = ''
            if r_idx < len(use_rows) and j < len(use_rows[r_idx]):
                val = use_rows[r_idx][j]
            tf = tbl.cell(i, j).text_frame
            set_paragraph_text(tf.paragraphs[0], val)
            for p in tf.paragraphs[1:]:
                clear_paragraph(p)


# -----------------------------
# Pictures
# -----------------------------

def get_blip_nodes(shape):
    return shape._element.xpath('.//*[local-name()="blip" or local-name()="svgBlip"]')


def related_part_by_rid(part, rid: str):
    return part.rels[rid].target_part


def swap_picture_embed(dest_slide, dest_shape_name: str, src_slide, src_shape_name: str):
    dest_shape = get_shape_by_name(dest_slide, dest_shape_name)
    src_shape = get_shape_by_name(src_slide, src_shape_name)

    src_blips = get_blip_nodes(src_shape)
    if not src_blips:
        raise ValueError(f"Source shape '{src_shape_name}' has no blip")

    src_rid = src_blips[0].get(qn('r:embed'))
    if not src_rid:
        raise ValueError(f"Source shape '{src_shape_name}' blip missing r:embed")

    img_part = related_part_by_rid(src_slide.part, src_rid)
    dest_part = dest_slide.part

    if hasattr(dest_part, 'relate_to'):
        new_rid = dest_part.relate_to(img_part, RT.IMAGE, is_external=False)
    else:
        new_rid = dest_part.rels.add_relationship(RT.IMAGE, img_part, None)

    for blip in get_blip_nodes(dest_shape):
        if blip.get(qn('r:embed')) is not None:
            blip.set(qn('r:embed'), new_rid)


# -----------------------------
# OPC slide duplication
# -----------------------------

def _xml(data: bytes):
    return ET.fromstring(data)


def _xml_bytes(elem) -> bytes:
    return ET.tostring(elem, xml_declaration=True, encoding='UTF-8', standalone='yes')


def duplicate_slides_in_package(template_path: str, out_path: str, duplicates: List[int], strip_notes_rels: bool = True) -> List[int]:
    if not duplicates:
        with open(template_path, 'rb') as fsrc:
            data = fsrc.read()
        with open(out_path, 'wb') as fdst:
            fdst.write(data)
        return []

    with zipfile.ZipFile(template_path, 'r') as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    slide_re = re.compile(r'^ppt/slides/slide(\d+)\.xml$')
    slide_nums = [int(slide_re.match(n).group(1)) for n in files if slide_re.match(n)]
    max_slide_num = max(slide_nums)

    rels_path = 'ppt/_rels/presentation.xml.rels'
    rels_xml = _xml(files[rels_path])
    rel_tag = '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'

    max_rid_num = 0
    for rel in rels_xml.findall(rel_tag):
        rid = rel.get('Id') or ''
        m = re.match(r'rId(\d+)$', rid)
        if m:
            max_rid_num = max(max_rid_num, int(m.group(1)))

    pres_path = 'ppt/presentation.xml'
    pres_xml = _xml(files[pres_path])
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = pres_xml.find(f'{{{p_ns}}}sldIdLst')
    if sldIdLst is None:
        raise RuntimeError('presentation.xml missing p:sldIdLst')

    max_sldid = 255
    for sldId in sldIdLst.findall(f'{{{p_ns}}}sldId'):
        try:
            max_sldid = max(max_sldid, int(sldId.get('id')))
        except Exception:
            pass

    ct_path = '[Content_Types].xml'
    ct_xml = _xml(files[ct_path])
    ct_ns = 'http://schemas.openxmlformats.org/package/2006/content-types'

    def add_ct_override(slide_num: int):
        part_name = f'/ppt/slides/slide{slide_num}.xml'
        for ov in ct_xml.findall(f'{{{ct_ns}}}Override'):
            if ov.get('PartName') == part_name:
                return
        ov = ET.Element(f'{{{ct_ns}}}Override')
        ov.set('PartName', part_name)
        ov.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
        ct_xml.append(ov)

    created: List[int] = []
    for src_slide_num in duplicates:
        max_slide_num += 1
        new_slide_num = max_slide_num

        src_xml_path = f'ppt/slides/slide{src_slide_num}.xml'
        dst_xml_path = f'ppt/slides/slide{new_slide_num}.xml'
        files[dst_xml_path] = files[src_xml_path]

        src_rels_path = f'ppt/slides/_rels/slide{src_slide_num}.xml.rels'
        dst_rels_path = f'ppt/slides/_rels/slide{new_slide_num}.xml.rels'
        if src_rels_path in files:
            rels = _xml(files[src_rels_path])
            if strip_notes_rels:
                for rel in list(rels.findall(rel_tag)):
                    if (rel.get('Type') or '').endswith('/notesSlide'):
                        rels.remove(rel)
            files[dst_rels_path] = _xml_bytes(rels)

        max_rid_num += 1
        new_rid = f'rId{max_rid_num}'
        rel = ET.Element(rel_tag)
        rel.set('Id', new_rid)
        rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
        rel.set('Target', f'slides/slide{new_slide_num}.xml')
        rels_xml.append(rel)

        max_sldid += 1
        sld = ET.Element(f'{{{p_ns}}}sldId')
        sld.set('id', str(max_sldid))
        sld.set(f'{{{r_ns}}}id', new_rid)
        sldIdLst.append(sld)

        add_ct_override(new_slide_num)
        created.append(new_slide_num)

    files[rels_path] = _xml_bytes(rels_xml)
    files[pres_path] = _xml_bytes(pres_xml)
    files[ct_path] = _xml_bytes(ct_xml)

    with zipfile.ZipFile(out_path, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)

    return created


def patch_docprops_app_xml(pptx_path: str):
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    pres = _xml(files['ppt/presentation.xml'])
    p_ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    slide_count = len(pres.findall('.//p:sldId', namespaces=p_ns))

    app_path = 'docProps/app.xml'
    if app_path not in files:
        return

    app = _xml(files[app_path])
    ext_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/extended-properties'
    slides_el = app.find(f'.//{{{ext_ns}}}Slides')
    if slides_el is not None:
        slides_el.text = str(slide_count)

    files[app_path] = _xml_bytes(app)

    tmp = pptx_path + '.tmp'
    with zipfile.ZipFile(tmp, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)
    os.replace(tmp, pptx_path)


# -----------------------------
# Generate from plan
# -----------------------------

def generate_from_plan(template_path: str, plan: Dict[str, Any], out_pptx: str, out_report: Optional[str], extra_warnings: Optional[List[Dict[str, Any]]] = None):
    deck = plan['deck']
    edits = plan['edits']

    source_slides = [int(item['source_slide']) for item in deck]

    counts: Dict[int, int] = {}
    duplicates: List[int] = []
    for s in source_slides:
        counts[s] = counts.get(s, 0) + 1
        if counts[s] > 1:
            duplicates.append(s)

    work_pptx = out_pptx + '.work.pptx'
    created = duplicate_slides_in_package(template_path, work_pptx, duplicates, strip_notes_rels=True)

    created_iter = iter(created)
    seen: Dict[int, int] = {}
    final_slide_numbers: List[int] = []
    for s in source_slides:
        seen[s] = seen.get(s, 0) + 1
        final_slide_numbers.append(s if seen[s] == 1 else next(created_iter))

    prs = Presentation(work_pptx)

    issues: List[Dict[str, Any]] = []
    if extra_warnings:
        issues.extend(extra_warnings)

    # family deletion
    for deck_pos, (item, slide_no) in enumerate(zip(deck, final_slide_numbers), start=1):
        sid = item['id']
        slide = prs.slides[slide_no - 1]
        entries = edits.get(sid, [])
        values_by_shape = {e['shape']: {'op': e['op'], 'value': e.get('value')} for e in entries}
        # Warnings: emit for any editable element left empty (LLM can decide if acceptable)
        for e in entries:
            op = e.get('op')
            if op == 'swap_picture':
                continue
            val = e.get('value')
            shp_name = e.get('shape')
            hint = (e.get('notes_hint') or '').strip()
            if is_empty_value(op, val):
                issues.append({
                    'severity': 'warn',
                    'slide': deck_pos,
                    'id': sid,
                    'shape': shp_name,
                    'issue': 'not_populated',
                    'detail': f"Shape value is empty. {hint}".strip(),
                })
            if op == 'set_table' and is_empty_value('set_table', val):
                issues.append({
                    'severity': 'warn',
                    'slide': deck_pos,
                    'id': sid,
                    'shape': shp_name,
                    'issue': 'table_not_populated',
                    'detail': f"Table has no headers/rows. {hint}".strip(),
                })

        prefixes, exacts, gap_warns = compute_family_delete_rules(values_by_shape)
        for w in gap_warns:
            issues.append({'severity':'warn','slide':deck_pos,'id':sid,'shape':f"{w['lead']}{w['empty_index']}_*",'issue':'non_tail_gap','detail':f"Index {w['empty_index']} is empty but later index {w['max_filled_index']} has content. Prefer removing from the end."})
        for pref in prefixes:
            for shp in list(iter_shapes(slide)):
                nm = getattr(shp, 'name', '') or ''
                if nm.startswith(pref):
                    try:
                        shp._element.getparent().remove(shp._element)
                    except Exception:
                        pass
        for exact in exacts:
            for shp in list(iter_shapes(slide)):
                nm = getattr(shp, 'name', '') or ''
                if nm == exact:
                    try:
                        shp._element.getparent().remove(shp._element)
                    except Exception:
                        pass

    # pictures first
    for deck_pos, (item, slide_no) in enumerate(zip(deck, final_slide_numbers), start=1):
        sid = item['id']
        slide = prs.slides[slide_no - 1]
        for e in edits.get(sid, []):
            if e.get('op') != 'swap_picture':
                continue
            asset = e.get('value')
            if is_empty_value('swap_picture', asset):
                continue
            src_slide = prs.slides[int(asset['asset_slide']) - 1]
            swap_picture_embed(slide, e['shape'], src_slide, asset['asset_shape'])

    # text + tables
    for deck_pos, (item, slide_no) in enumerate(zip(deck, final_slide_numbers), start=1):
        sid = item['id']
        slide = prs.slides[slide_no - 1]
        for e in edits.get(sid, []):
            op = e.get('op')
            if op == 'swap_picture':
                continue
            try:
                shp = get_shape_by_name(slide, e['shape'])
            except KeyError:
                continue
            val = e.get('value')
            if op == 'set_text':
                set_shape_text_single_paragraph(shp, val or '')
            elif op == 'set_paragraphs':
                set_multi_paragraph(shp, val or [])
            elif op == 'set_bullets':
                set_bullets_preserve_style(shp, val or [])
            elif op == 'set_table':
                v = val or {}
                set_table_safe(slide, e['shape'], v.get('headers', []) or [], v.get('rows', []) or [])

    # reorder/prune
    sldIdLst = prs.slides._sldIdLst
    elems = list(sldIdLst)
    keep = set(final_slide_numbers)
    rids_to_drop = [sldId.get(qn('r:id')) for idx, sldId in enumerate(elems, start=1) if idx not in keep]

    new_elems = [elems[i - 1] for i in final_slide_numbers]
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)
    for el in new_elems:
        sldIdLst.append(el)

    for rid in rids_to_drop:
        try:
            prs.part.drop_rel(rid)
        except Exception:
            pass

    prs.save(out_pptx)
    patch_docprops_app_xml(out_pptx)

    if out_report:
        with open(out_report, 'w', encoding='utf-8') as f:
            json.dump(issues, f, indent=2, ensure_ascii=False)
    else:
        print(json.dumps(issues, indent=2, ensure_ascii=False))

    try:
        os.remove(work_pptx)
    except Exception:
        pass


# -----------------------------
# Generate from request
# -----------------------------

def generate_from_request(cli_template_path: str, request_path: str, out_pptx: str, out_report: Optional[str], emit_expanded_plan: Optional[str] = None):
    req = json.load(open(request_path, 'r', encoding='utf-8'))

    req_template = req.get('template')
    if cli_template_path == '-':
        if not (isinstance(req_template, str) and req_template.strip() and os.path.exists(req_template)):
            raise SystemExit("request.template must exist when using '-' as template argument")
        template = req_template
    else:
        template = req_template if (isinstance(req_template, str) and os.path.exists(req_template)) else cli_template_path

    prs = Presentation(template)
    plan, warnings, report = expand_request_to_plan(prs, template, req)

    illegal = report.get('illegal')
    if illegal:
        illegal_ids = sorted(set(x.get('id') for x in illegal if x.get('id')))
        allowed_all = report.get('allowed', {})
        allowed_filtered = {sid: allowed_all.get(sid) for sid in illegal_ids if sid in allowed_all}
        payload = {
            'status': 'error',
            'message': 'Illegal targets in request',
            'template_used': template,
            'illegal': illegal,
            'allowed_by_slide': allowed_filtered,
            'warnings': warnings,
        }
        print(json.dumps(payload, indent=2, ensure_ascii=False), file=sys.stderr)
        if out_report:
            json.dump(payload, open(out_report, 'w', encoding='utf-8'), indent=2, ensure_ascii=False)
        raise SystemExit(2)

    if emit_expanded_plan:
        json.dump(plan, open(emit_expanded_plan, 'w', encoding='utf-8'), indent=2, ensure_ascii=False)

    generate_from_plan(template, plan, out_pptx, out_report, extra_warnings=warnings)


# -----------------------------
# CLI
# -----------------------------

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('mode', choices=['extract', 'generate'])
    ap.add_argument('pptx', help="Template path or '-' to force request.template")
    ap.add_argument('--request', default=None)
    ap.add_argument('-o', '--out', default=None)
    ap.add_argument('--report', default=None)
    ap.add_argument('--emit-expanded-plan', default=None)

    args = ap.parse_args()

    if args.mode == 'extract':
        md = extract_notes_markdown(args.pptx)
        if args.out:
            open(args.out, 'w', encoding='utf-8').write(md)
        else:
            print(md)
        return

    if args.mode == 'generate':
        if not args.request:
            raise SystemExit('generate requires --request')
        out = args.out or 'out.pptx'
        generate_from_request(args.pptx, args.request, out, args.report, emit_expanded_plan=args.emit_expanded_plan)
        print('Saved:', out)
        return


if __name__ == '__main__':
    main()
