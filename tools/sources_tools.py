"""
Sources and Deliverables Tools
Allows the agent to inspect what the user has dropped into sources/,
read text-based briefs, and save finished presentations to deliverables/.
"""

import base64
import io
import json
import os
import re
import shutil
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any

from ibm_watsonx_orchestrate.agent_builder.tools import tool

# ── Path resolution ────────────────────────────────────────────────────────────
# Works both locally (relative to project root) and inside the WXO runtime
# (where __file__ is inside the bundled package directory).
_TOOL_DIR   = Path(__file__).resolve().parent          # …/presentation-system/tools
_SYSTEM_DIR = _TOOL_DIR.parent                         # …/presentation-system
_PROJECT_DIR = _SYSTEM_DIR.parent                      # …/adk-project

def _resolve(folder: str) -> Path:
    """Return an absolute Path for a project-level folder, trying several roots."""
    for base in (_PROJECT_DIR, _SYSTEM_DIR, Path.cwd()):
        candidate = base / folder
        if candidate.exists():
            return candidate
    # If it doesn't exist yet, create it relative to the project root
    path = _PROJECT_DIR / folder
    path.mkdir(parents=True, exist_ok=True)
    return path

SOURCES_DIR     = _resolve("sources")
DELIVERABLES_DIR = _resolve("deliverables")

# ── Helpers ────────────────────────────────────────────────────────────────────

_TEXT_EXTS  = {".md", ".txt", ".csv", ".json", ".yaml", ".yml"}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".webp"}
_DECK_EXTS  = {".pptx", ".ppt", ".potx"}
_DATA_EXTS  = {".xlsx", ".xls", ".csv"}

def _categorise(suffix: str) -> str:
    s = suffix.lower()
    if s in _IMAGE_EXTS:  return "image"
    if s in _DECK_EXTS:   return "deck"
    if s in _DATA_EXTS:   return "data"
    if s in _TEXT_EXTS:   return "text"
    return "other"

def _safe_filename(name: str) -> str:
    """Sanitise a string for use as a filename."""
    name = re.sub(r'[^\w\s\-]', '', name)
    name = re.sub(r'\s+', '_', name.strip())
    return name[:80] or "presentation"


# ══════════════════════════════════════════════════════════════════════════════
#  Tool 1 — list_sources
# ══════════════════════════════════════════════════════════════════════════════

@tool(
    name="list_sources",
    display_name="List Source Files",
    description=(
        "List all files currently in the sources/ directory. "
        "Call this to see what logos, briefs, data files, and reference materials "
        "the user has made available before building the presentation."
    ),
)
def list_sources() -> str:
    """Return metadata for every file in sources/.

    Returns:
        JSON string with keys:
          - files: list of {name, category, size_kb, modified, is_logo}
          - counts: breakdown by category
          - logos: names of detected logo files (filename contains 'logo')
          - text_briefs: names of readable text files
          - empty: bool — True if no files found
    """
    files = []
    for p in sorted(SOURCES_DIR.iterdir()):
        if p.name.startswith(".") or p.name == "README.md" or not p.is_file():
            continue
        stat = p.stat()
        category = _categorise(p.suffix)
        files.append({
            "name":     p.name,
            "category": category,
            "size_kb":  round(stat.st_size / 1024, 1),
            "modified": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M"),
            "is_logo":  "logo" in p.name.lower(),
        })

    counts: Dict[str, int] = {}
    for f in files:
        counts[f["category"]] = counts.get(f["category"], 0) + 1

    return json.dumps({
        "files":       files,
        "counts":      counts,
        "logos":       [f["name"] for f in files if f["is_logo"]],
        "text_briefs": [f["name"] for f in files if f["category"] == "text"],
        "empty":       len(files) == 0,
    }, indent=2)


# ══════════════════════════════════════════════════════════════════════════════
#  Tool 2 — read_source_text
# ══════════════════════════════════════════════════════════════════════════════

@tool(
    name="read_source_text",
    display_name="Read Source Text File",
    description=(
        "Read the text content of a file in sources/. "
        "Supports .txt, .md, .csv, .json, .yaml files. "
        "Use this to ingest briefs, talking points, data tables, or any written "
        "content the user has dropped in before building the presentation."
    ),
)
def read_source_text(filename: str) -> str:
    """Read and return the text content of a file in sources/.

    Args:
        filename: Name of the file in sources/ (e.g. 'brief.md', 'metrics.csv').

    Returns:
        JSON string with keys:
          - filename, content (str), lines (int), size_kb (float)
        Or an error key if the file cannot be read.
    """
    path = SOURCES_DIR / os.path.basename(filename)

    if not path.exists():
        return json.dumps({"error": f"File '{filename}' not found in sources/"})

    if path.suffix.lower() not in _TEXT_EXTS:
        return json.dumps({
            "error": (
                f"'{filename}' is not a readable text file "
                f"(type: {_categorise(path.suffix)}). "
                "Use list_sources to see what's available."
            )
        })

    try:
        content = path.read_text(encoding="utf-8", errors="replace")
        return json.dumps({
            "filename": path.name,
            "content":  content,
            "lines":    content.count("\n") + 1,
            "size_kb":  round(path.stat().st_size / 1024, 1),
        }, ensure_ascii=False)
    except Exception as exc:
        return json.dumps({"error": f"Could not read '{filename}': {exc}"})


# ══════════════════════════════════════════════════════════════════════════════
#  Tool 3 — save_deliverable
# ══════════════════════════════════════════════════════════════════════════════

@tool(
    name="save_deliverable",
    display_name="Save Presentation to Deliverables",
    description=(
        "Save a generated PPTX (returned by create_presentation) to the "
        "deliverables/ folder with a human-readable filename. "
        "Optionally injects a client or IBM logo from sources/ onto the title slide. "
        "Returns the saved filename and path."
    ),
)
def save_deliverable(
    pptx_bytes_b64: str,
    presentation_title: str,
    logo_filename: Optional[str] = None,
) -> str:
    """Save a base64-encoded PPTX to deliverables/ and optionally inject a logo.

    Args:
        pptx_bytes_b64: Base64-encoded PPTX bytes (the output of create_presentation
                        encoded with base64.b64encode(...).decode()).
        presentation_title: Human-readable title used for the filename,
                            e.g. "Q4 Results Review" → Q4_Results_Review.pptx
        logo_filename: Optional filename of a logo image in sources/ to inject
                       into the bottom-right of the title slide
                       (e.g. 'client_logo.png'). Leave empty to skip.

    Returns:
        JSON string with keys:
          - saved_path: relative path to the saved file
          - filename: just the filename
          - size_kb: file size
          - logo_injected: bool — whether a logo was successfully added
          - logo_error: str — reason if logo injection failed (empty if ok)
    """
    # Decode PPTX bytes
    try:
        pptx_bytes = base64.b64decode(pptx_bytes_b64)
    except Exception as exc:
        return json.dumps({"error": f"Invalid base64 PPTX data: {exc}"})

    # Build output path
    safe_title = _safe_filename(presentation_title)
    filename   = f"{safe_title}.pptx"
    out_path   = DELIVERABLES_DIR / filename

    # Avoid overwriting: append timestamp if file exists
    if out_path.exists():
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename  = f"{safe_title}_{ts}.pptx"
        out_path  = DELIVERABLES_DIR / filename

    # Write PPTX to disk
    out_path.write_bytes(pptx_bytes)

    logo_injected = False
    logo_error    = ""

    # Optionally inject logo
    if logo_filename:
        logo_path = SOURCES_DIR / os.path.basename(logo_filename)
        if not logo_path.exists():
            logo_error = f"Logo file '{logo_filename}' not found in sources/"
        elif logo_path.suffix.lower() not in _IMAGE_EXTS:
            logo_error = f"'{logo_filename}' is not a recognised image file"
        else:
            try:
                from pptx import Presentation as _Prs
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN

                prs = _Prs(str(out_path))
                slide = prs.slides[0]  # title slide

                # Place logo in bottom-right corner (1.5" × 0.6" @ right edge)
                slide_w = prs.slide_width
                slide_h = prs.slide_height
                logo_w  = Inches(1.5)
                logo_h  = Inches(0.6)
                left    = slide_w  - logo_w  - Inches(0.2)
                top     = slide_h  - logo_h  - Inches(0.15)

                slide.shapes.add_picture(str(logo_path), left, top, logo_w, logo_h)
                prs.save(str(out_path))
                logo_injected = True

            except Exception as exc:
                logo_error = f"Logo injection failed: {exc}"

    return json.dumps({
        "saved_path":    f"deliverables/{filename}",
        "filename":      filename,
        "size_kb":       round(out_path.stat().st_size / 1024, 1),
        "logo_injected": logo_injected,
        "logo_error":    logo_error,
    })

# Made with Bob


# ══════════════════════════════════════════════════════════════════════════════
#  Tool 7 — read_source_data
# ══════════════════════════════════════════════════════════════════════════════

@tool(
    name="read_source_data",
    display_name="Read Source Data File",
    description=(
        "Read structured data from an Excel (.xlsx, .xls) or CSV (.csv) file in sources/. "
        "Returns all sheets with headers, rows, and a plain-English summary of what was found. "
        "Use this before building a presentation when the user has dropped in metrics, "
        "KPIs, financial data, survey results, or any tabular data they want on slides."
    ),
)
def read_source_data(filename: str, max_rows: Optional[int] = 100) -> str:
    """Read and return structured data from a .xlsx, .xls, or .csv file in sources/.

    Args:
        filename: Name of the data file in sources/ (e.g. 'metrics.xlsx', 'kpis.csv').
        max_rows: Maximum rows to return per sheet (default 100). Use a smaller value
                  for large files to keep the response concise.

    Returns:
        JSON string with keys:
          - filename, file_type
          - sheets: list of {sheet_name, headers, rows, row_count, numeric_columns,
                              summary (plain-English description of what the data contains)}
          - total_sheets, total_rows
          - suggested_slides: list of suggested slide types based on the data found
          - error: set if the file could not be read
    """
    path = SOURCES_DIR / os.path.basename(filename)

    if not path.exists():
        return json.dumps({"error": f"File '{filename}' not found in sources/"})

    suffix = path.suffix.lower()
    if suffix not in {".xlsx", ".xls", ".csv"}:
        return json.dumps({
            "error": (
                f"'{filename}' is not a supported data file. "
                "Supported formats: .xlsx, .xls, .csv"
            )
        })

    sheets_out = []
    total_rows = 0

    try:
        if suffix == ".csv":
            # ── CSV ──────────────────────────────────────────────────────────
            import csv
            with open(path, newline="", encoding="utf-8-sig", errors="replace") as f:
                reader = csv.reader(f)
                all_rows = list(reader)

            if not all_rows:
                return json.dumps({"error": "CSV file is empty"})

            headers  = [str(h).strip() for h in all_rows[0]]
            data_rows = [
                [str(cell).strip() for cell in row]
                for row in all_rows[1 : (max_rows or 100) + 1]
            ]
            total_rows = len(all_rows) - 1

            numeric_cols = _detect_numeric_columns(headers, data_rows)
            summary      = _summarise_sheet("Sheet1", headers, data_rows, numeric_cols)

            sheets_out.append({
                "sheet_name":      "Sheet1",
                "headers":         headers,
                "rows":            data_rows,
                "row_count":       total_rows,
                "numeric_columns": numeric_cols,
                "summary":         summary,
            })

        else:
            # ── Excel ────────────────────────────────────────────────────────
            try:
                import openpyxl
            except ImportError:
                return json.dumps({"error": "openpyxl is required to read .xlsx files. Add it to requirements.txt."})

            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            cap = max_rows or 100

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                all_rows = list(ws.iter_rows(values_only=True))
                if not all_rows:
                    continue

                # First non-empty row = headers
                headers = [str(h).strip() if h is not None else "" for h in all_rows[0]]
                data_rows = []
                for row in all_rows[1 : cap + 1]:
                    data_rows.append([
                        str(cell).strip() if cell is not None else ""
                        for cell in row
                    ])
                sheet_total = len(all_rows) - 1
                total_rows += sheet_total

                numeric_cols = _detect_numeric_columns(headers, data_rows)
                summary      = _summarise_sheet(sheet_name, headers, data_rows, numeric_cols)

                sheets_out.append({
                    "sheet_name":      sheet_name,
                    "headers":         headers,
                    "rows":            data_rows,
                    "row_count":       sheet_total,
                    "numeric_columns": numeric_cols,
                    "summary":         summary,
                })

            wb.close()

        if not sheets_out:
            return json.dumps({"error": "No data found in file"})

        suggested = _suggest_slides(sheets_out)

        return json.dumps({
            "filename":        path.name,
            "file_type":       suffix.lstrip("."),
            "sheets":          sheets_out,
            "total_sheets":    len(sheets_out),
            "total_rows":      total_rows,
            "suggested_slides": suggested,
        }, ensure_ascii=False, indent=2)

    except Exception as exc:
        return json.dumps({"error": f"Could not read '{filename}': {exc}"})


# ── Data reading helpers ───────────────────────────────────────────────────────

def _detect_numeric_columns(headers: list, rows: list) -> list:
    """Return list of column names that appear to contain numeric values."""
    numeric = []
    for col_idx, header in enumerate(headers):
        values = [row[col_idx] for row in rows if col_idx < len(row) and row[col_idx]]
        if not values:
            continue
        numeric_count = 0
        for v in values[:20]:
            cleaned = v.replace(",", "").replace("%", "").replace("£", "").replace("$", "").replace("€", "").strip()
            try:
                float(cleaned)
                numeric_count += 1
            except ValueError:
                pass
        if numeric_count / len(values[:20]) >= 0.7:
            numeric.append(header)
    return numeric


def _summarise_sheet(sheet_name: str, headers: list, rows: list, numeric_cols: list) -> str:
    """Generate a plain-English summary of what a sheet contains."""
    if not rows:
        return f"Sheet '{sheet_name}' has headers but no data rows."

    parts = [f"{len(rows)} rows × {len(headers)} columns."]

    if numeric_cols:
        parts.append(f"Numeric columns: {', '.join(numeric_cols[:6])}.")

        # Pull out standout values from numeric columns
        highlights = []
        for col_name in numeric_cols[:4]:
            if col_name not in headers:
                continue
            col_idx = headers.index(col_name)
            values = []
            for row in rows:
                if col_idx < len(row) and row[col_idx]:
                    cleaned = row[col_idx].replace(",", "").replace("%", "").strip()
                    try:
                        values.append(float(cleaned))
                    except ValueError:
                        pass
            if values:
                highlights.append(f"{col_name}: min={min(values):.1f}, max={max(values):.1f}")
        if highlights:
            parts.append("Key ranges — " + "; ".join(highlights) + ".")

    text_cols = [h for h in headers if h and h not in numeric_cols]
    if text_cols:
        parts.append(f"Category columns: {', '.join(text_cols[:4])}.")

    return " ".join(parts)


def _suggest_slides(sheets: list) -> list:
    """Suggest slide types based on what data was found."""
    suggestions = []
    for sheet in sheets:
        numeric = sheet.get("numeric_columns", [])
        rows    = sheet.get("rows", [])
        name    = sheet.get("sheet_name", "data")
        count   = len(numeric)

        if count >= 2:
            suggestions.append({
                "slide_type": "stats",
                "reason": f"Sheet '{name}' has {count} numeric columns — good for a stats slide with 2–4 hero numbers.",
                "columns_to_use": numeric[:4],
            })
        if len(rows) >= 3 and len(sheet.get("headers", [])) >= 2:
            suggestions.append({
                "slide_type": "table",
                "reason": f"Sheet '{name}' has {len(rows)} rows — could work as a table slide.",
            })
        if len(rows) <= 2 and count >= 1:
            suggestions.append({
                "slide_type": "content",
                "reason": f"Sheet '{name}' has summary figures — use as bullet points on a content slide.",
            })

    if not suggestions:
        suggestions.append({
            "slide_type": "content",
            "reason": "Use the data as bullet points on a content or two_column slide.",
        })

    return suggestions


# ══════════════════════════════════════════════════════════════════════════════
#  Tool 4 — list_deliverables
# ══════════════════════════════════════════════════════════════════════════════

@tool(
    name="list_deliverables",
    display_name="List Deliverable Presentations",
    description=(
        "List all presentations saved in the deliverables/ folder, newest first. "
        "Shows filename, size, and creation date. "
        "Use this to see what has already been generated, pick a file to revise, "
        "or show the user what is ready to download."
    ),
)
def list_deliverables() -> str:
    """Return metadata for every .pptx in deliverables/, newest first.

    Returns:
        JSON string with keys:
          - files: list of {name, size_kb, modified, is_named (bool — has readable name)}
          - count: total number of files
          - named: files with readable names (not UUID-only)
          - latest: name of the most recently modified file, or null
    """
    files = []
    for p in DELIVERABLES_DIR.iterdir():
        if not p.is_file() or p.suffix.lower() != ".pptx":
            continue
        if p.name.startswith("~$"):   # skip Office lock files
            continue
        stat  = p.stat()
        is_named = not re.match(
            r'^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}\.pptx$',
            p.name, re.I
        )
        files.append({
            "name":     p.name,
            "size_kb":  round(stat.st_size / 1024, 1),
            "modified": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M"),
            "is_named": is_named,
        })

    files.sort(key=lambda f: f["modified"], reverse=True)

    return json.dumps({
        "files":  files,
        "count":  len(files),
        "named":  [f["name"] for f in files if f["is_named"]],
        "latest": files[0]["name"] if files else None,
    }, indent=2)


# ══════════════════════════════════════════════════════════════════════════════
#  Tool 5 — delete_file
# ══════════════════════════════════════════════════════════════════════════════

@tool(
    name="delete_file",
    display_name="Delete Source or Deliverable File",
    description=(
        "Delete a file from sources/ or deliverables/. "
        "Specify the folder ('sources' or 'deliverables') and the filename. "
        "Use with caution — deletion is permanent."
    ),
)
def delete_file(folder: str, filename: str) -> str:
    """Delete a file from sources/ or deliverables/.

    Args:
        folder: Either 'sources' or 'deliverables'.
        filename: Name of the file to delete (basename only, no path).

    Returns:
        JSON with keys: deleted (bool), path, error (if any).
    """
    if folder not in ("sources", "deliverables"):
        return json.dumps({"deleted": False, "error": "folder must be 'sources' or 'deliverables'"})

    base = SOURCES_DIR if folder == "sources" else DELIVERABLES_DIR
    path = base / os.path.basename(filename)

    if not path.exists():
        return json.dumps({"deleted": False, "error": f"'{filename}' not found in {folder}/"})

    if path.name == "README.md":
        return json.dumps({"deleted": False, "error": "Cannot delete README.md"})

    try:
        path.unlink()
        return json.dumps({"deleted": True, "path": f"{folder}/{filename}"})
    except Exception as exc:
        return json.dumps({"deleted": False, "error": str(exc)})


# ══════════════════════════════════════════════════════════════════════════════
#  Tool 6 — rename_deliverable
# ══════════════════════════════════════════════════════════════════════════════

@tool(
    name="rename_deliverable",
    display_name="Rename Deliverable",
    description=(
        "Rename a file in deliverables/ to a human-readable name. "
        "Useful for giving a UUID-named file a proper title, "
        "e.g. '3d734bc2.pptx' → 'Q4_Results_Review.pptx'."
    ),
)
def rename_deliverable(current_filename: str, new_title: str) -> str:
    """Rename a deliverable to a human-readable filename.

    Args:
        current_filename: Current filename in deliverables/ (e.g. '3d734bc2-....pptx').
        new_title: Human-readable title for the new filename
                   (e.g. 'Q4 Results Review' → Q4_Results_Review.pptx).

    Returns:
        JSON with keys: renamed (bool), old_name, new_name, error (if any).
    """
    src = DELIVERABLES_DIR / os.path.basename(current_filename)
    if not src.exists():
        return json.dumps({"renamed": False, "error": f"'{current_filename}' not found in deliverables/"})

    new_name = _safe_filename(new_title) + ".pptx"
    dst = DELIVERABLES_DIR / new_name

    if dst.exists():
        ts      = datetime.now().strftime("%Y%m%d_%H%M%S")
        new_name = _safe_filename(new_title) + f"_{ts}.pptx"
        dst      = DELIVERABLES_DIR / new_name

    try:
        src.rename(dst)
        return json.dumps({"renamed": True, "old_name": src.name, "new_name": new_name})
    except Exception as exc:
        return json.dumps({"renamed": False, "error": str(exc)})
